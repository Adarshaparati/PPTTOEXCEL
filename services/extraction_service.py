import os
import base64
import hashlib
from io import BytesIO
from pptx import Presentation
from openpyxl import Workbook
from typing import Optional, Dict, List
from datetime import datetime
from services.s3_service import s3_service

try:
    from PIL import Image
except ImportError:
    Image = None

class ExtractionService:
    """Service for extracting data from PowerPoint presentations"""
    
    @staticmethod
    def sanitize_text(text):
        """Remove illegal characters for Excel cells"""
        if text is None:
            return ''
        if not isinstance(text, str):
            text = str(text)
        import re
        text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', text)
        if len(text) > 32000:
            text = text[:32000] + '...'
        return text
    
    @staticmethod
    def emu_to_inches(emu):
        """Convert EMU (English Metric Units) to inches"""
        return round(emu / 914400, 2) if emu else 0
    
    @staticmethod
    def get_color_info(color_obj):
        """Extract color information"""
        try:
            if hasattr(color_obj, 'rgb'):
                return f"RGB({color_obj.rgb.red},{color_obj.rgb.green},{color_obj.rgb.blue})"
            elif hasattr(color_obj, 'theme_color'):
                return f"Theme Color {color_obj.theme_color}"
            else:
                return "Auto/Default"
        except:
            return "Unknown"
    
    @staticmethod
    def get_image_extension(content_type):
        """Get file extension from content type"""
        extension_map = {
            'image/jpeg': '.jpg',
            'image/jpg': '.jpg',
            'image/png': '.png',
            'image/gif': '.gif',
            'image/bmp': '.bmp',
            'image/tiff': '.tiff',
            'image/webp': '.webp',
            'image/svg+xml': '.svg'
        }
        return extension_map.get(content_type.lower(), '.img')
    
    @staticmethod
    def get_mime_type(content_type):
        """Get proper MIME type for base64 data URL"""
        if content_type and content_type.startswith('image/'):
            return content_type
        return 'image/jpeg'
    
    def extract_font_info(self, text_frame):
        """Extract comprehensive font information from text frame"""
        font_info = {
            'name': 'Default',
            'size': 'Default',
            'bold': False,
            'italic': False,
            'underline': False,
            'color': 'Auto',
            'alignment': 'Left',
            'line_spacing': 'Default',
            'paragraph_spacing': 'Default'
        }
        
        try:
            if text_frame and text_frame.paragraphs:
                first_paragraph = text_frame.paragraphs[0]
                
                if hasattr(first_paragraph, 'alignment'):
                    alignment_map = {
                        0: 'Left', 1: 'Center', 2: 'Right', 3: 'Justify',
                        4: 'Distribute', 5: 'Thai Distribute'
                    }
                    font_info['alignment'] = alignment_map.get(first_paragraph.alignment, 'Left')
                
                if hasattr(first_paragraph, 'line_spacing'):
                    try:
                        if first_paragraph.line_spacing:
                            font_info['line_spacing'] = f"{first_paragraph.line_spacing:.2f}"
                    except:
                        pass
                
                try:
                    space_before = first_paragraph.space_before.pt if first_paragraph.space_before else 0
                    space_after = first_paragraph.space_after.pt if first_paragraph.space_after else 0
                    font_info['paragraph_spacing'] = f"Before:{space_before}pt, After:{space_after}pt"
                except:
                    pass
                
                if first_paragraph.runs:
                    first_run = first_paragraph.runs[0]
                    font = first_run.font
                    
                    font_info['name'] = font.name or 'Default'
                    font_info['size'] = f"{font.size.pt}pt" if font.size else 'Default'
                    font_info['bold'] = font.bold if font.bold is not None else False
                    font_info['italic'] = font.italic if font.italic is not None else False
                    font_info['underline'] = font.underline if font.underline is not None else False
                    font_info['color'] = self.get_color_info(font.color) if hasattr(font, 'color') else 'Auto'
        except:
            pass
        
        return font_info
    
    def extract_fill_info(self, shape):
        """Extract detailed fill information"""
        fill_info = {
            'color': 'None',
            'type': 'None',
            'transparency': 0
        }
        
        try:
            if hasattr(shape, 'fill') and shape.fill:
                fill_type_map = {
                    0: 'No Fill', 1: 'Solid', 2: 'Gradient', 3: 'Picture',
                    4: 'Pattern', 5: 'Group', 6: 'Background'
                }
                fill_info['type'] = fill_type_map.get(shape.fill.type, 'Unknown')
                
                if hasattr(shape.fill, 'fore_color'):
                    fill_info['color'] = self.get_color_info(shape.fill.fore_color)
                
                if hasattr(shape.fill, 'transparency'):
                    fill_info['transparency'] = f"{shape.fill.transparency * 100:.1f}%"
        except:
            pass
        
        return fill_info
    
    def extract_image_info(self, shape, slide_num, shape_index, upload_to_s3=True):
        """Extract detailed image information and optionally upload to S3"""
        image_info = {
            'has_image': False,
            'format': '',
            'width': 0,
            'height': 0,
            'file_size': 0,
            'url': '',
            'base64': '',
            's3_url': ''
        }
        
        try:
            if shape.shape_type == 12 or hasattr(shape, 'image'):
                image_info['has_image'] = True
                
                if hasattr(shape, 'image') and shape.image:
                    image_blob = shape.image.blob
                    image_info['format'] = shape.image.content_type or 'Unknown'
                    image_info['file_size'] = len(image_blob) if image_blob else 0
                    
                    if image_blob:
                        # Generate unique filename
                        image_hash = hashlib.md5(image_blob).hexdigest()[:12]
                        file_extension = self.get_image_extension(image_info['format'])
                        filename = f"slide_{slide_num}_shape_{shape_index}_{image_hash}{file_extension}"
                        
                        # Save locally
                        images_folder = "extracted_images"
                        if not os.path.exists(images_folder):
                            os.makedirs(images_folder)
                        
                        filepath = os.path.join(images_folder, filename)
                        with open(filepath, 'wb') as f:
                            f.write(image_blob)
                        
                        image_info['url'] = filepath.replace('\\', '/')
                        
                        # Upload to S3 if enabled
                        if upload_to_s3:
                            s3_result = s3_service.upload_image_to_s3(
                                image_data=image_blob,
                                filename=filename,
                                content_type=image_info['format']
                            )
                            if s3_result.get('success'):
                                image_info['s3_url'] = s3_result.get('s3_url', '')
                        
                        # Create base64 data URL
                        base64_string = base64.b64encode(image_blob).decode('utf-8')
                        mime_type = self.get_mime_type(image_info['format'])
                        image_info['base64'] = f"data:{mime_type};base64,{base64_string}"
                        
                        # Get image dimensions
                        try:
                            if Image and image_blob:
                                img = Image.open(BytesIO(image_blob))
                                image_info['width'] = img.width
                                image_info['height'] = img.height
                            else:
                                image_info['width'] = self.emu_to_inches(shape.width)
                                image_info['height'] = self.emu_to_inches(shape.height)
                        except:
                            image_info['width'] = self.emu_to_inches(shape.width)
                            image_info['height'] = self.emu_to_inches(shape.height)
        except:
            pass
        
        return image_info
    
    def extract_chart_info(self, shape):
        """Extract detailed chart information"""
        chart_info = {
            'type': 'None',
            'title': 'None',
            'data': 'None',
            'categories': 'None',
            'series': 'None'
        }
        
        try:
            if shape.shape_type == 3 and hasattr(shape, 'chart'):
                chart = shape.chart
                
                chart_type_map = {
                    1: 'Area', 2: 'Bar', 3: 'Column', 4: 'Line', 5: 'Pie',
                    6: 'Scatter', 7: 'Surface', 8: 'Radar', 9: 'Treemap',
                    10: 'Sunburst', 11: 'Histogram', 12: 'BoxWhisker',
                    13: 'Waterfall', 14: 'Funnel', 15: 'Map'
                }
                chart_info['type'] = chart_type_map.get(chart.chart_type, f'Unknown({chart.chart_type})')
                
                if hasattr(chart, 'chart_title') and chart.chart_title:
                    try:
                        if hasattr(chart.chart_title, 'text_frame') and chart.chart_title.text_frame:
                            chart_info['title'] = chart.chart_title.text_frame.text.strip()
                    except:
                        chart_info['title'] = 'Has Title'
                
                series_data = []
                categories_data = []
                chart_data_points = []
                
                try:
                    if hasattr(chart, 'plots') and chart.plots:
                        plot = chart.plots[0]
                        
                        if hasattr(plot, 'categories') and plot.categories:
                            try:
                                categories_data = [str(cat) for cat in plot.categories]
                            except:
                                categories_data = ['Category data available']
                        
                        if hasattr(plot, 'series'):
                            for i, series in enumerate(plot.series):
                                series_name = getattr(series, 'name', f'Series {i+1}')
                                series_data.append(series_name)
                                
                                try:
                                    if hasattr(series, 'values') and series.values:
                                        values = [str(v) for v in series.values if v is not None]
                                        if values:
                                            chart_data_points.append(f"{series_name}: [{', '.join(values[:10])}{'...' if len(values) > 10 else ''}]")
                                except:
                                    chart_data_points.append(f"{series_name}: [Values available]")
                except:
                    pass
                
                if categories_data:
                    chart_info['categories'] = ' | '.join(categories_data[:10]) + ('...' if len(categories_data) > 10 else '')
                
                if series_data:
                    chart_info['series'] = ' | '.join(series_data)
                
                if chart_data_points:
                    chart_info['data'] = ' | '.join(chart_data_points)
                
                if chart_info['data'] == 'None':
                    try:
                        chart_xml = shape.element.xml
                        if 'val' in chart_xml and 'cat' in chart_xml:
                            chart_info['data'] = '[Chart data embedded in XML]'
                    except:
                        pass
        except Exception as e:
            if shape.shape_type == 3:
                chart_info['type'] = 'Chart (extraction failed)'
                chart_info['data'] = f'Chart present but data extraction failed: {str(e)[:50]}'
        
        return chart_info
    
    def extract_ppt_to_excel(self, ppt_bytes: bytes, upload_images_to_s3: bool = True) -> tuple:
        """
        Extract PowerPoint data to Excel
        
        Args:
            ppt_bytes: PowerPoint file content as bytes
            upload_images_to_s3: Whether to upload extracted images to S3
            
        Returns:
            tuple: (workbook, extracted_data_dict)
        """
        prs = Presentation(BytesIO(ppt_bytes))
        wb = Workbook()
        ws = wb.active
        ws.title = "PPT_Data"
        
        headers = [
            "Slide No", "Shape Name", "Shape Type", "Content",
            "Left (EMU)", "Top (EMU)", "Width (EMU)", "Height (EMU)",
            "Left (Inches)", "Top (Inches)", "Width (Inches)", "Height (Inches)",
            "Font Name", "Font Size", "Font Bold", "Font Italic", "Font Underline", "Font Color",
            "Text Alignment", "Line Spacing", "Paragraph Spacing",
            "Fill Color", "Fill Type", "Transparency", "Line Color", "Line Width", "Line Style", "Rotation",
            "Has Image", "Image Format", "Image Width", "Image Height", "Image File Size", "Image URL", "Image S3 URL", "Image Base64",
            "Chart Type", "Chart Title", "Chart Data", "Chart Categories", "Chart Series",
            "Hyperlink", "Z-Order", "Hidden", "Shadow", "Glow Effect", "Reflection",
            "3D Effects", "Placeholder Type", "Animation Effects"
        ]
        ws.append(headers)
        
        extracted_data = {
            'total_slides': len(prs.slides),
            'slides': [],
            'total_images': 0,
            'total_charts': 0,
            'total_tables': 0
        }
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_data = {
                'slide_number': slide_num,
                'shapes': []
            }
            
            for shape_index, shape in enumerate(slide.shapes):
                content = ""
                hyperlink = ""
                
                if shape.has_text_frame:
                    content = self.sanitize_text(shape.text.strip().replace("\n", " | "))
                    try:
                        if shape.text_frame.paragraphs:
                            for para in shape.text_frame.paragraphs:
                                for run in para.runs:
                                    if hasattr(run, 'hyperlink') and run.hyperlink.address:
                                        hyperlink = self.sanitize_text(run.hyperlink.address)
                                        break
                    except:
                        pass
                
                elif shape.shape_type == 19:  # Table
                    table_data = []
                    try:
                        for row in shape.table.rows:
                            row_text = [self.sanitize_text(cell.text.strip()) for cell in row.cells]
                            table_data.append(", ".join(row_text))
                        content = " | ".join(table_data)
                        extracted_data['total_tables'] += 1
                    except:
                        content = "[TABLE - Could not extract data]"
                
                elif shape.shape_type == 12:  # Picture
                    content = "[IMAGE]"
                    extracted_data['total_images'] += 1
                
                elif shape.shape_type == 3:  # Chart
                    content = "[CHART]"
                    extracted_data['total_charts'] += 1
                
                # Extract all information
                font_info = self.extract_font_info(shape.text_frame if shape.has_text_frame else None)
                fill_info = self.extract_fill_info(shape)
                image_info = self.extract_image_info(shape, slide_num, shape_index, upload_images_to_s3)
                chart_info = self.extract_chart_info(shape)
                
                # Position and size
                left_emu = shape.left
                top_emu = shape.top
                width_emu = shape.width
                height_emu = shape.height
                
                row_data = [
                    slide_num, self.sanitize_text(shape.name), self.get_shape_type_name(shape.shape_type), content,
                    left_emu, top_emu, width_emu, height_emu,
                    self.emu_to_inches(left_emu), self.emu_to_inches(top_emu),
                    self.emu_to_inches(width_emu), self.emu_to_inches(height_emu),
                    self.sanitize_text(font_info['name']), self.sanitize_text(font_info['size']),
                    font_info['bold'], font_info['italic'], font_info['underline'],
                    self.sanitize_text(font_info['color']), self.sanitize_text(font_info['alignment']),
                    self.sanitize_text(font_info['line_spacing']), self.sanitize_text(font_info['paragraph_spacing']),
                    self.sanitize_text(fill_info['color']), self.sanitize_text(fill_info['type']),
                    self.sanitize_text(fill_info['transparency']), "None", "None", "None", 0,
                    image_info['has_image'], self.sanitize_text(image_info['format']),
                    image_info['width'], image_info['height'], image_info['file_size'],
                    self.sanitize_text(image_info['url']), self.sanitize_text(image_info['s3_url']),
                    self.sanitize_text(image_info['base64']),
                    self.sanitize_text(chart_info['type']), self.sanitize_text(chart_info['title']),
                    self.sanitize_text(chart_info['data']), self.sanitize_text(chart_info['categories']),
                    self.sanitize_text(chart_info['series']),
                    hyperlink, shape_index, False, False, False, False, False, "Not a placeholder", "None"
                ]
                
                ws.append(row_data)
                
                # Add to JSON data with ALL fields matching Excel
                shape_data = {
                    'slide_number': slide_num,
                    'shape_name': self.sanitize_text(shape.name),
                    'shape_type': self.get_shape_type_name(shape.shape_type),
                    'content': content,
                    'left_emu': left_emu,
                    'top_emu': top_emu,
                    'width_emu': width_emu,
                    'height_emu': height_emu,
                    'left_inches': self.emu_to_inches(left_emu),
                    'top_inches': self.emu_to_inches(top_emu),
                    'width_inches': self.emu_to_inches(width_emu),
                    'height_inches': self.emu_to_inches(height_emu),
                    'font_name': self.sanitize_text(font_info['name']),
                    'font_size': self.sanitize_text(font_info['size']),
                    'font_bold': font_info['bold'],
                    'font_italic': font_info['italic'],
                    'font_underline': font_info['underline'],
                    'font_color': self.sanitize_text(font_info['color']),
                    'text_alignment': self.sanitize_text(font_info['alignment']),
                    'line_spacing': self.sanitize_text(font_info['line_spacing']),
                    'paragraph_spacing': self.sanitize_text(font_info['paragraph_spacing']),
                    'fill_color': self.sanitize_text(fill_info['color']),
                    'fill_type': self.sanitize_text(fill_info['type']),
                    'transparency': self.sanitize_text(fill_info['transparency']),
                    'line_color': "None",
                    'line_width': "None",
                    'line_style': "None",
                    'rotation': 0,
                    'has_image': image_info['has_image'],
                    'image_format': self.sanitize_text(image_info['format']),
                    'image_width': image_info['width'],
                    'image_height': image_info['height'],
                    'image_file_size': image_info['file_size'],
                    'image_url': self.sanitize_text(image_info['url']),
                    'image_s3_url': self.sanitize_text(image_info['s3_url']),
                    'image_base64': self.sanitize_text(image_info['base64']),
                    'chart_type': self.sanitize_text(chart_info['type']),
                    'chart_title': self.sanitize_text(chart_info['title']),
                    'chart_data': self.sanitize_text(chart_info['data']),
                    'chart_categories': self.sanitize_text(chart_info['categories']),
                    'chart_series': self.sanitize_text(chart_info['series']),
                    'hyperlink': hyperlink,
                    'z_order': shape_index,
                    'hidden': False,
                    'shadow': False,
                    'glow_effect': False,
                    'reflection': False,
                    '3d_effects': False,
                    'placeholder_type': "Not a placeholder",
                    'animation_effects': "None"
                }
                slide_data['shapes'].append(shape_data)
            
            extracted_data['slides'].append(slide_data)
        
        return wb, extracted_data
    
    @staticmethod
    def get_shape_type_name(shape_type):
        """Convert shape type number to readable name"""
        type_mapping = {
            1: "AutoShape", 2: "Callout", 3: "Chart", 4: "Comment",
            5: "Freeform", 6: "Group", 7: "Line", 8: "LinkedOLEObject",
            9: "LinkedPicture", 10: "Media", 11: "OLEObject", 12: "Picture",
            13: "Placeholder", 14: "TextBox", 15: "3DModel", 16: "Canvas",
            17: "Connector", 18: "Ink", 19: "Table", 20: "SmartArt"
        }
        return type_mapping.get(shape_type, f"Unknown({shape_type})")

# Create singleton instance
extraction_service = ExtractionService()
