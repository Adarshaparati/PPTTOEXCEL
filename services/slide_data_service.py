import os
import io
import requests
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Any, Optional
from PIL import Image
from services.s3_service import s3_service
from datetime import datetime

# Import slide handlers
from services.handlers.points import handle_points_slide
from services.handlers.image_text import handle_image_text_slide
from services.handlers.table import handle_table_slide
from services.handlers.phases import handle_phases_slide
from services.handlers.statistics import handle_statistics_slide
from services.handlers.people import handle_people_slide
from services.handlers.cover import handle_cover_slide
from services.handlers.contact import handle_contact_slide
from services.handlers.images import handle_images_slide
from services.handlers.graphs import handle_graph_slide


class SlideDataService:
    """Service for managing and generating different slide types"""
    
    @staticmethod
    def hex_to_rgb(hex_str):
        """Convert #RRGGBB string to RGBColor safely."""
        if not hex_str:
            return RGBColor(0, 0, 0)
        
        hex_str = str(hex_str).strip().replace("#", "")
        if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
            print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
            return RGBColor(0, 0, 0)
        
        return RGBColor(int(hex_str[0:2], 16),
                        int(hex_str[2:4], 16),
                        int(hex_str[4:6], 16))
    
    def download_image_from_url(self, image_url: str) -> BytesIO:
        """
        Download image from URL and return as BytesIO
        
        Args:
            image_url: URL of the image to download
            
        Returns:
            BytesIO object containing the image data
        """
        try:
            print(f"🖼️ Downloading image from: {image_url}")
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            print(f"✅ Image downloaded successfully ({len(response.content)} bytes)")
            return BytesIO(response.content)
        except Exception as e:
            error_msg = f"Failed to download image from {image_url}: {str(e)}"
            print(f"❌ {error_msg}")
            raise RuntimeError(error_msg)
    
    def download_template_from_s3(self, s3_url: str) -> BytesIO:
        """
        Download PowerPoint template from S3
        
        Args:
            s3_url: S3 URL of the template or just the S3 key
            
        Returns:
            BytesIO object containing the template data
        """
        try:
            from urllib.parse import unquote, urlparse
            
            # Extract key from S3 URL
            if s3_url.startswith('http'):
                # Try direct HTTP download first (works for public S3 URLs)
                print(f"📥 Attempting direct download from S3 URL: {s3_url}")
                try:
                    response = requests.get(s3_url, timeout=30)
                    response.raise_for_status()
                    print(f"✅ Template downloaded directly via HTTP ({len(response.content)} bytes)")
                    return BytesIO(response.content)
                except Exception as http_error:
                    print(f"⚠️ Direct HTTP download failed: {http_error}")
                    print(f"🔄 Falling back to boto3 download...")
                    
                    # Fall back to boto3 download
                    parsed = urlparse(s3_url)
                    # Get the path without leading slash and decode any URL encoding
                    key = unquote(parsed.path.lstrip('/'))
                    
                    if not key:
                        raise ValueError(f"Invalid S3 URL format - no key found: {s3_url}")
                        
                    print(f"📥 Parsed S3 URL - Key: {key}")
            else:
                # Assume it's just the key
                key = s3_url
                print(f"📥 Using provided key: {key}")
            
            # Download from S3 using boto3
            file_data = s3_service.download_file(key)

            if not file_data:
                raise RuntimeError(f"S3 download failed or returned no data for key: {key}")

            print(f"✅ Template downloaded successfully ({len(file_data)} bytes)")
            return BytesIO(file_data)
            
        except Exception as e:
            error_msg = f"Failed to download template from S3: {str(e)}"
            print(f"❌ {error_msg}")
            raise RuntimeError(error_msg)
    
    def generate_slide(
        self,
        slide_type: str,
        presentation_s3_url: str,
        slide_data: Dict[str, Any]
    ) -> BytesIO:
        """
        Universal slide generator - routes to appropriate handler based on slide type
        
        Args:
            slide_type: Type of slide ('points', 'image_text', 'table', 'phases', etc.)
            presentation_s3_url: S3 URL of the presentation to modify
            slide_data: Dictionary containing slide configuration
            
        Returns:
            BytesIO object containing the modified presentation
        """
        print(f"� Generating {slide_type.upper()} slide...")
        
        # Download the presentation
        presentation_bytes = self.download_template_from_s3(presentation_s3_url)
        
        # Download image if provided
        image_data = None
        if slide_data.get('image_url'):
            try:
                image_data = self.download_image_from_url(slide_data['image_url'])
            except Exception as e:
                print(f"⚠️ Warning: Could not download image: {e}")
        
        # Route to appropriate handler
        handler_map = {
            'points': handle_points_slide,
            'image_text': handle_image_text_slide,
            'table': handle_table_slide,
            'phases': handle_phases_slide,
            'statistics': handle_statistics_slide,
            'people': handle_people_slide,
            'cover': handle_cover_slide,
            'contact': handle_contact_slide,
            'images': handle_images_slide,
            'graphs': handle_graph_slide,
        }
        
        handler = handler_map.get(slide_type)
        if not handler:
            available_types = ', '.join(handler_map.keys())
            raise ValueError(f"Unsupported slide type '{slide_type}'. Available types: {available_types}")
        
        # Call the appropriate handler
        try:
            result = handler(presentation_bytes, slide_data, image_data)
            print(f"✅ {slide_type.upper()} slide generated successfully")
            return result
        except Exception as e:
            error_msg = f"Error in {slide_type} handler: {str(e)}"
            print(f"❌ {error_msg}")
            raise RuntimeError(error_msg)
    
    # Convenience methods for backward compatibility and specific slide types
    def generate_points_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate points slide - backwards compatible method"""
        return self.generate_slide('points', template_s3_url, slide_data)
    
    def generate_image_text_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate image+text slide - backwards compatible method"""
        return self.generate_slide('image_text', template_s3_url, slide_data)
    
    def generate_table_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate table slide - backwards compatible method"""
        return self.generate_slide('table', template_s3_url, slide_data)
    
    def generate_phases_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate phases slide - backwards compatible method"""
        return self.generate_slide('phases', template_s3_url, slide_data)
    
    def generate_statistics_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate statistics slide - backwards compatible method"""
        return self.generate_slide('statistics', template_s3_url, slide_data)
    
    def generate_people_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate people slide - backwards compatible method"""
        return self.generate_slide('people', template_s3_url, slide_data)
    
    def generate_cover_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate cover slide - backwards compatible method"""
        return self.generate_slide('cover', template_s3_url, slide_data)
    
    def generate_contact_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate contact slide - backwards compatible method"""
        return self.generate_slide('contact', template_s3_url, slide_data)
    
    def generate_images_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate images slide - backwards compatible method"""
        return self.generate_slide('images', template_s3_url, slide_data)
    
    def generate_graph_slide(self, template_s3_url: str, slide_data: Dict[str, Any]) -> BytesIO:
        """Generate graph slide - backwards compatible method"""
        return self.generate_slide('graphs', template_s3_url, slide_data)
    
    def generate_multi_slide_presentation(
        self,
        template_s3_url: str,
        slides_config: List[Dict[str, Any]]
    ) -> BytesIO:
        """
        OPTIMIZED: Generate multiple slides in a single presentation
        
        Downloads template ONCE, loads presentation ONCE, modifies all slides,
        then saves ONCE for maximum efficiency.
        
        Args:
            template_s3_url: S3 URL of the PowerPoint template
            slides_config: List of slide configurations, each containing:
                - slide_type: str - Type of slide ('points', 'image_text', 'table', 'phases', etc.)
                - slide_data: Dict - Data specific to the slide type
                
        Returns:
            BytesIO object containing the complete presentation
        """
        print(f"🎨 Generating multi-slide presentation (OPTIMIZED)...")
        
        # Download template ONCE from S3
        presentation_bytes = self.download_template_from_s3(template_s3_url)
        
        # Load presentation ONCE into memory
        prs = Presentation(presentation_bytes)
        print(f"📄 Loaded presentation with {len(prs.slides)} slides")
        
        # Process each slide configuration
        slides_processed = 0
        for idx, config in enumerate(slides_config, 1):
            slide_type = config.get('slide_type')
            slide_data = config.get('slide_data', {})
            
            if not slide_type:
                print(f"⚠️ Skipping config {idx} without slide_type")
                continue
            
            try:
                print(f"🔧 Processing slide {idx}/{len(slides_config)}: {slide_type}")
                self._modify_slide_in_presentation(prs, slide_type, slide_data)
                slides_processed += 1
            except Exception as e:
                print(f"⚠️ Error processing {slide_type} slide: {e}")
                continue
        
        # Save presentation ONCE to BytesIO
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        print(f"✅ Multi-slide presentation generated successfully ({slides_processed}/{len(slides_config)} slides processed)")
        return output
    
    def _modify_slide_in_presentation(self, prs: Presentation, slide_type: str, slide_data: Dict[str, Any]) -> None:
        """
        OPTIMIZED: Modify a single slide directly in the presentation object
        
        This method works directly on the Presentation object instead of using
        BytesIO chaining, which is more efficient for multi-slide generation.
        
        Args:
            prs: Presentation object to modify
            slide_type: Type of slide ('points', 'image_text', 'table', 'phases', etc.)
            slide_data: Dictionary with slide configuration
        """
        from io import BytesIO
        
        # Get slide number and validate
        slide_number = slide_data.get('slide_number', 1)
        slide_index = slide_number - 1
        
        if slide_index >= len(prs.slides) or slide_index < 0:
            raise ValueError(f"Slide {slide_number} not found in presentation (has {len(prs.slides)} slides)")
        
        slide = prs.slides[slide_index]
        print(f"📝 Modifying slide {slide_number} ({slide_type})...")
        
        # Download image if needed
        image_data = None
        if slide_data.get('image_url'):
            try:
                image_data = self.download_image_from_url(slide_data['image_url'])
            except Exception as e:
                print(f"⚠️ Warning: Could not download image for {slide_type}: {e}")
        
        # Handle multiple images for cover/contact/images slide types
        if slide_data.get('image') and isinstance(slide_data['image'], list):
            # These slide types handle multiple images internally
            pass
        
        # Route to appropriate modification logic based on slide type
        if slide_type == 'points':
            self._modify_points_slide(slide, slide_data, image_data)
        elif slide_type == 'image_text':
            self._modify_image_text_slide(slide, slide_data, image_data)
        elif slide_type == 'table':
            self._modify_table_slide(slide, slide_data)
        elif slide_type == 'phases':
            self._modify_phases_slide(slide, slide_data, image_data)
        elif slide_type == 'statistics':
            self._modify_statistics_slide(slide, slide_data, image_data)
        elif slide_type == 'people':
            self._modify_people_slide(slide, slide_data, image_data)
        elif slide_type == 'cover':
            self._modify_cover_slide(slide, slide_data)
        elif slide_type == 'contact':
            self._modify_contact_slide(slide, slide_data)
        elif slide_type == 'images':
            self._modify_images_slide(slide, slide_data)
        elif slide_type == 'graphs':
            self._modify_graph_slide(slide, slide_data)
        else:
            raise ValueError(f"Unsupported slide type: {slide_type}")
        
        print(f"✅ Slide {slide_number} modified successfully")
    
    # Import helper functions from handlers to reuse logic
    def _modify_points_slide(self, slide, slide_data: Dict[str, Any], image_data: BytesIO = None):
        """Modify points slide directly"""
        from services.handlers.points import _update_header, _update_description, _update_image, _update_points
        
        if slide_data.get('header'):
            _update_header(slide, slide_data)
        if slide_data.get('description'):
            _update_description(slide, slide_data)
        if image_data and slide_data.get('image_url'):
            _update_image(slide, image_data)
        if slide_data.get('points'):
            _update_points(slide, slide_data)
    
    def _modify_image_text_slide(self, slide, slide_data: Dict[str, Any], image_data: BytesIO = None):
        """Modify image+text slide directly"""
        from services.handlers.image_text import _update_title, _update_text, _update_image
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('text'):
            _update_text(slide, slide_data)
        if image_data and slide_data.get('image_url'):
            _update_image(slide, image_data)
    
    def _modify_table_slide(self, slide, slide_data: Dict[str, Any]):
        """Modify table slide directly"""
        from services.handlers.table import _update_title, _update_table
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('table_data'):
            _update_table(slide, slide_data)
    
    def _modify_phases_slide(self, slide, slide_data: Dict[str, Any], image_data: BytesIO = None):
        """Modify phases slide directly"""
        from services.handlers.phases import _update_title, _update_phases, _update_image
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('phases'):
            _update_phases(slide, slide_data)
        if image_data and slide_data.get('image_url'):
            _update_image(slide, image_data)
    
    def _modify_statistics_slide(self, slide, slide_data: Dict[str, Any], image_data: BytesIO = None):
        """Modify statistics slide directly"""
        from services.handlers.statistics import _update_title, _update_description, _update_statistics, _update_background_color, _update_image
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('description'):
            _update_description(slide, slide_data)
        if slide_data.get('stat_data'):
            _update_statistics(slide, slide_data)
        if slide_data.get('background_color'):
            _update_background_color(slide, slide_data['background_color'])
        if image_data and slide_data.get('image_url'):
            _update_image(slide, image_data)
    
    def _modify_people_slide(self, slide, slide_data: Dict[str, Any], image_data: BytesIO = None):
        """Modify people slide directly"""
        from services.handlers.people import _update_title, _update_description, _update_people, _update_background_color, _update_image
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('description'):
            _update_description(slide, slide_data)
        if slide_data.get('names'):
            _update_people(slide, slide_data)
        if slide_data.get('background_color'):
            _update_background_color(slide, slide_data['background_color'])
        if image_data and slide_data.get('image_url'):
            _update_image(slide, image_data)
    
    def _modify_cover_slide(self, slide, slide_data: Dict[str, Any]):
        """Modify cover slide directly"""
        from services.handlers.cover import _update_title, _update_subtitle, _update_company_name, _update_images_from_urls, _apply_color_scheme
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('subtitle'):
            _update_subtitle(slide, slide_data)
        if slide_data.get('company_name'):
            _update_company_name(slide, slide_data)
        if slide_data.get('image'):
            _update_images_from_urls(slide, slide_data.get('image', []))
        if slide_data.get('colors'):
            _apply_color_scheme(slide, slide_data.get('colors', {}))
    
    def _modify_contact_slide(self, slide, slide_data: Dict[str, Any]):
        """Modify contact slide directly"""
        from services.handlers.contact import _update_title, _update_website, _update_linkedin, _update_email, _update_phone, _update_images_from_urls, _apply_color_scheme
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('website_link'):
            _update_website(slide, slide_data)
        if slide_data.get('linkedin_link'):
            _update_linkedin(slide, slide_data)
        if slide_data.get('contact_email'):
            _update_email(slide, slide_data)
        if slide_data.get('contact_phone'):
            _update_phone(slide, slide_data)
        if slide_data.get('image'):
            _update_images_from_urls(slide, slide_data.get('image', []))
        if slide_data.get('colors'):
            _apply_color_scheme(slide, slide_data.get('colors', {}))
    
    def _modify_images_slide(self, slide, slide_data: Dict[str, Any]):
        """Modify images slide directly"""
        from services.handlers.images import _update_title, _update_image_gallery
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('images'):
            _update_image_gallery(slide, slide_data)
    
    def _modify_graph_slide(self, slide, slide_data: Dict[str, Any]):
        """Modify graph slide directly"""
        from services.handlers.graphs import _update_title, _update_chart
        
        if slide_data.get('title'):
            _update_title(slide, slide_data)
        if slide_data.get('chart_data'):
            _update_chart(slide, slide_data)
    
    def _process_slide_in_presentation(self, presentation_bytes: BytesIO, slide_type: str, slide_data: Dict[str, Any]) -> BytesIO:
        """
        LEGACY: Process a single slide in a presentation (BytesIO chaining)
        
        This method is kept for backward compatibility but is less efficient
        than the new _modify_slide_in_presentation method used in multi-slide generation.
        """
        # Download image if needed
        image_data = None
        if slide_data.get('image_url'):
            try:
                image_data = self.download_image_from_url(slide_data['image_url'])
            except Exception as e:
                print(f"⚠️ Warning: Could not download image for {slide_type}: {e}")
        
        # Route to appropriate handler
        handler_map = {
            'points': handle_points_slide,
            'image_text': handle_image_text_slide,
            'table': handle_table_slide,
            'phases': handle_phases_slide,
            'statistics': handle_statistics_slide,
            'people': handle_people_slide,
            'cover': handle_cover_slide,
            'contact': handle_contact_slide,
            'images': handle_images_slide,
            'graphs': handle_graph_slide,
        }
        
        handler = handler_map.get(slide_type)
        if not handler:
            print(f"⚠️ Unknown slide type: {slide_type}")
            return presentation_bytes
        
        return handler(presentation_bytes, slide_data, image_data)


# Create singleton instance
slide_data_service = SlideDataService()