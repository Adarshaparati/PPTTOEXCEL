"""
Points slide handler - Handles generation/modification of bullet point slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_points_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle points/bullet slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
        image_data: Optional BytesIO containing image data
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"🎯 Processing POINTS slide...")
    
    prs = Presentation(presentation_bytes)
    
    # Get slide
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update header
    if slide_data.get('header'):
        _update_header(slide, slide_data)
    
    # Update description
    if slide_data.get('description'):
        _update_description(slide, slide_data)
    
    # Update image
    if image_data and slide_data.get('image_url'):
        _update_image(slide, image_data)
    
    # Update points/bullets
    if slide_data.get('points'):
        _update_points(slide, slide_data)
    
    # Save to BytesIO
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Points slide processed successfully")
    return output


def _update_header(slide, slide_data: Dict[str, Any]):
    """Update header text and formatting"""
    for shape in slide.shapes:
        if shape.name == 'Header1' or (hasattr(shape, 'text_frame') and 'header' in shape.name.lower()):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['header']
                if slide_data.get('header_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['header_color'])
                print(f"📝 Updated header: {slide_data['header']}")
                break


def _update_description(slide, slide_data: Dict[str, Any]):
    """Update description text and formatting"""
    for shape in slide.shapes:
        if shape.name == 'Description1' or (hasattr(shape, 'text_frame') and 'description' in shape.name.lower()):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['description']
                if slide_data.get('description_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['description_color'])
                print(f"📝 Updated description: {slide_data['description'][:50]}...")
                break


def _update_image(slide, image_data: BytesIO):
    """Update image in slide"""
    for shape in slide.shapes:
        if shape.name == 'Image' or 'image' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                # Get position and size
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                # Remove old image
                sp = shape.element
                sp.getparent().remove(sp)
                # Add new image
                slide.shapes.add_picture(image_data, left, top, width, height)
                print(f"🖼️ Updated image")
                break


def _update_points(slide, slide_data: Dict[str, Any]):
    """Update bullet points"""
    points_text = '\n'.join([f"• {point.get('text', '')}" for point in slide_data['points']])
    
    for shape in slide.shapes:
        if shape.name.startswith('Description1_BG') or 'points' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = points_text
                # Apply point-specific formatting if provided
                for idx, (paragraph, point) in enumerate(zip(shape.text_frame.paragraphs, slide_data['points'])):
                    if point.get('color'):
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(point['color'])
                    if point.get('font_size'):
                        for run in paragraph.runs:
                            run.font.size = Pt(point['font_size'])
                print(f"📋 Updated {len(slide_data['points'])} bullet points")
                break