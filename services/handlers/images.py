"""
Images slide handler - Handles generation/modification of multi-image gallery slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any
import requests


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_images_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle images/gallery slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - slide_name: str - Name of the slide
            - title: str - Slide title
            - headers: List[str] - Array of headers for each image
            - descriptions: List[str] - Array of descriptions for each image
            - images: List[str] - Array of image URLs
        image_data: Optional BytesIO containing image data (not used directly, images come from URLs)
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"🖼️ Processing IMAGES slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update images and their metadata
    if slide_data.get('images'):
        _update_image_gallery(slide, slide_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Images slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_image_gallery(slide, slide_data: Dict[str, Any]):
    """Update image gallery with headers and descriptions"""
    images = slide_data.get('images', [])
    headers = slide_data.get('headers', [])
    descriptions = slide_data.get('descriptions', [])
    
    if not images:
        return
    
    # Ensure all arrays have the same length
    max_length = len(images)
    headers = headers + [''] * (max_length - len(headers))
    descriptions = descriptions + [''] * (max_length - len(descriptions))
    
    # Process each image
    for i, (image_url, header, description) in enumerate(zip(images, headers, descriptions), 1):
        # Update image
        _update_image(slide, i, image_url)
        
        # Update header
        if header:
            _update_image_header(slide, i, header)
        
        # Update description
        if description:
            _update_image_description(slide, i, description)


def _update_image(slide, index: int, image_url: str):
    """Update a specific image in the gallery"""
    image_shape_name = f"Image{index}"
    
    for shape in slide.shapes:
        if shape.name == image_shape_name or f"image{index}" in shape.name.lower() or f"picture{index}" in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                try:
                    print(f"🖼️ Downloading image {index} from: {image_url}")
                    response = requests.get(image_url, timeout=10)
                    response.raise_for_status()
                    
                    image_data = BytesIO(response.content)
                    
                    # Get position and size
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    
                    # Remove old image
                    sp = shape.element
                    sp.getparent().remove(sp)
                    
                    # Add new image
                    slide.shapes.add_picture(image_data, left, top, width, height)
                    print(f"✅ Updated image {index}")
                    
                except Exception as e:
                    print(f"⚠️ Failed to update image {index} from {image_url}: {e}")
                break


def _update_image_header(slide, index: int, header: str):
    """Update header for a specific image"""
    header_shape_name = f"Header{index}"
    
    for shape in slide.shapes:
        if shape.name == header_shape_name or f"header{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = header
                
                # Make header bold
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                
                print(f"📝 Updated header {index}: {header}")
                break


def _update_image_description(slide, index: int, description: str):
    """Update description for a specific image"""
    description_shape_name = f"Description{index}"
    
    for shape in slide.shapes:
        if shape.name == description_shape_name or f"description{index}" in shape.name.lower() or f"desc{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = description
                print(f"📝 Updated description {index}: {description[:50]}...")
                break
