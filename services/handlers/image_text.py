"""
Image+Text slide handler - Handles generation/modification of image and text slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from typing import Dict, Any


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_image_text_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle image+text slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
        image_data: Optional BytesIO containing image data
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"🖼️ Processing IMAGE+TEXT slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update text
    if slide_data.get('text'):
        _update_text(slide, slide_data)
    
    # Update image
    if image_data and slide_data.get('image_url'):
        _update_image(slide, image_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Image+Text slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if shape.name == 'P100' or 'title' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                if slide_data.get('title_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['title_color'])
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_text(slide, slide_data: Dict[str, Any]):
    """Update text content and formatting"""
    for shape in slide.shapes:
        if shape.name == 'S100' or 'text' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['text']
                if slide_data.get('text_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['text_color'])
                print(f"📝 Updated text: {slide_data['text'][:50]}...")
                break


def _update_image(slide, image_data: BytesIO):
    """Update image in slide"""
    for shape in slide.shapes:
        if 'image' in shape.name.lower() or 'picture' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(image_data, left, top, width, height)
                print(f"🖼️ Updated image")
                break