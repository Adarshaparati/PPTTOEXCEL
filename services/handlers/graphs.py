"""
Graph slide handler - Handles generation/modification of chart/graph slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any, List
import json


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_graph_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle graph/chart slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - title: str - Slide title
            - chart_type: str - Type of chart (bar, line, pie, etc.)
            - charts: List[Dict] - List of charts with title and chart_data
            - chart_options: Dict - Chart configuration options
        image_data: Optional BytesIO containing image data (not used for graphs)
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"📈 Processing GRAPH slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update charts
    if slide_data.get('charts'):
        _update_charts(slide, slide_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Graph slide processed successfully")
    return output


def handle_graph_slide_optimized(slide, slide_data: Dict[str, Any]):
    """
    OPTIMIZED: Handle graph/chart slide directly on slide object.
    Used by multi-slide generation for better performance.
    
    Args:
        slide: PowerPoint slide object
        slide_data: Dictionary containing graph slide data
    """
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update charts
    if slide_data.get('charts'):
        _update_charts(slide, slide_data)


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_charts(slide, slide_data: Dict[str, Any]):
    """Update chart data and formatting"""
    charts = slide_data.get('charts', [])
    chart_type = slide_data.get('chart_type', 'bar')
    chart_options = slide_data.get('chart_options', {})
    
    if not charts:
        return
    
    # Process each chart
    for chart_idx, chart_config in enumerate(charts, 1):
        _update_single_chart(slide, chart_idx, chart_config, chart_type, chart_options)


def _update_single_chart(slide, chart_index: int, chart_config: Dict[str, Any], chart_type: str, chart_options: Dict[str, Any]):
    """Update a single chart in the slide"""
    chart_title = chart_config.get('title', '')
    chart_data = chart_config.get('chart_data', {})
    
    # Look for chart shapes (Chart1, Chart2, etc.) or text areas to store chart data
    chart_shape_name = f"Chart{chart_index}"
    chart_title_name = f"ChartTitle{chart_index}"
    chart_data_name = f"ChartData{chart_index}"
    
    # Update chart title
    for shape in slide.shapes:
        if (shape.name == chart_title_name or 
            f"charttitle{chart_index}" in shape.name.lower() or
            (chart_index == 1 and 'chart' in shape.name.lower() and 'title' in shape.name.lower())):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = chart_title
                print(f"📊 Updated chart {chart_index} title: {chart_title}")
                break
    
    # Update chart data (stored as text or processed by template)
    chart_info = _format_chart_data_as_text(chart_data, chart_type)
    
    for shape in slide.shapes:
        if (shape.name == chart_data_name or 
            f"chartdata{chart_index}" in shape.name.lower() or
            f"chart{chart_index}" in shape.name.lower()):
            
            # If it's a text shape, store formatted chart data
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = chart_info
                print(f"📊 Updated chart {chart_index} data")
                break
            
            # If it's an actual chart object (PowerPoint chart)
            elif hasattr(shape, 'chart'):
                try:
                    _update_powerpoint_chart(shape.chart, chart_data, chart_type)
                    print(f"📊 Updated PowerPoint chart {chart_index}")
                    break
                except Exception as e:
                    print(f"⚠️ Could not update PowerPoint chart {chart_index}: {e}")
                    # Fallback: try to find a text shape to store data
                    continue


def _format_chart_data_as_text(chart_data: Dict[str, Any], chart_type: str) -> str:
    """Format chart data as readable text for text shapes"""
    labels = chart_data.get('labels', [])
    datasets = chart_data.get('datasets', [])
    
    if not labels or not datasets:
        return f"Chart Type: {chart_type.title()}\nNo data available"
    
    lines = [f"Chart Type: {chart_type.title()}"]
    lines.append(f"Labels: {', '.join(str(label) for label in labels)}")
    
    for dataset in datasets:
        dataset_label = dataset.get('label', 'Series')
        dataset_data = dataset.get('data', [])
        lines.append(f"{dataset_label}: {', '.join(str(value) for value in dataset_data)}")
    
    return '\n'.join(lines)


def _update_powerpoint_chart(chart, chart_data: Dict[str, Any], chart_type: str):
    """Update actual PowerPoint chart object"""
    try:
        # This is a simplified implementation for PowerPoint chart objects
        # The exact implementation depends on the chart structure in your template
        
        labels = chart_data.get('labels', [])
        datasets = chart_data.get('datasets', [])
        
        if not labels or not datasets:
            return
        
        # Access chart data
        chart_data_obj = chart.chart_data
        
        # Clear existing data (if possible)
        # Note: PowerPoint chart manipulation is complex and may vary by chart type
        
        # Update categories (labels)
        categories = chart_data_obj.categories
        for i, label in enumerate(labels[:len(categories)]):
            if i < len(categories):
                categories[i] = str(label)
        
        # Update series data
        for series_idx, dataset in enumerate(datasets[:len(chart.series)]):
            if series_idx < len(chart.series):
                series = chart.series[series_idx]
                series.name = dataset.get('label', f'Series {series_idx + 1}')
                
                data_values = dataset.get('data', [])
                for i, value in enumerate(data_values[:len(series.values)]):
                    if i < len(series.values):
                        try:
                            series.values[i] = float(value) if str(value).replace('.', '').isdigit() else 0
                        except (ValueError, TypeError):
                            series.values[i] = 0
        
        print(f"📊 Updated PowerPoint chart with {len(labels)} categories and {len(datasets)} series")
        
    except Exception as e:
        print(f"⚠️ PowerPoint chart update failed: {e}")
        raise


def _apply_chart_options(slide, chart_options: Dict[str, Any]):
    """Apply chart styling options"""
    if not chart_options:
        return
    
    background_color = chart_options.get('backgroundColor')
    if background_color:
        # Apply background color to chart areas or slide background
        for shape in slide.shapes:
            if 'chart' in shape.name.lower() or 'background' in shape.name.lower():
                # This would need template-specific implementation
                print(f"🎨 Chart background color: {background_color}")