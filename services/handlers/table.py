"""
Table slide handler - Handles generation/modification of table slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from typing import Dict, Any, List


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_table_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle table slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
        image_data: Optional BytesIO containing image data (not used for tables)
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"📊 Processing TABLE slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update table
    if slide_data.get('table_data'):
        _update_table(slide, slide_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Table slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_table(slide, slide_data: Dict[str, Any]):
    """Update table data and formatting"""
    table_data = slide_data.get('table_data', [])
    if not table_data:
        return
    
    for shape in slide.shapes:
        if shape.has_table:
            table = shape.table
            
            # Update table cells
            rows_updated = 0
            for row_idx, row_data in enumerate(table_data):
                if row_idx >= len(table.rows):
                    break
                for col_idx, cell_data in enumerate(row_data):
                    if col_idx >= len(table.columns):
                        break
                    cell = table.cell(row_idx, col_idx)
                    cell.text = str(cell_data)
                    
                    # Apply header formatting
                    if row_idx == 0 and slide_data.get('header_row', True):
                        if slide_data.get('header_color'):
                            for paragraph in cell.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = hex_to_rgb(slide_data['header_color'])
                                    run.font.bold = True
                rows_updated += 1
            
            print(f"📋 Updated table with {rows_updated} rows and {len(table_data[0]) if table_data else 0} columns")
            break