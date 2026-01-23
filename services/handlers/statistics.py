"""
Statistics slide handler - Handles generation/modification of statistics slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_statistics_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle statistics slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - title: str - Slide title
            - description: str - Optional description
            - stat_data: List[Dict] - List of statistics with label, value, color, font_size
            - title_color: str - Color for title
            - description_color: str - Color for description
            - background_color: str - Color for background
        image_data: Optional BytesIO containing image data
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"📊 Processing STATISTICS slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update description
    if slide_data.get('description'):
        _update_description(slide, slide_data)
    
    # Update statistics
    if slide_data.get('stat_data'):
        _update_statistics(slide, slide_data)
    
    # Update background color if provided
    if slide_data.get('background_color'):
        _update_background_color(slide, slide_data['background_color'])
    
    # Update image if provided
    if image_data and slide_data.get('image_url'):
        _update_image(slide, image_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Statistics slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                if slide_data.get('title_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['title_color'])
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_description(slide, slide_data: Dict[str, Any]):
    """Update description text and formatting"""
    for shape in slide.shapes:
        if 'description' in shape.name.lower() or shape.name.startswith('Description'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['description']
                if slide_data.get('description_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['description_color'])
                print(f"📝 Updated description: {slide_data['description'][:50]}...")
                break


def _update_statistics(slide, slide_data: Dict[str, Any]):
    """Update statistics data"""
    stat_data = slide_data.get('stat_data', [])
    if not stat_data:
        return
    
    # Look for stat shapes (Stat1, Stat2, etc.) or (Label1/Value1, Label2/Value2, etc.)
    for i, stat in enumerate(stat_data, 1):
        label = stat.get('label', '')
        value = stat.get('value', '')
        
        # Try to find combined stat shape first
        stat_shape_name = f"Stat{i}"
        found = False
        
        for shape in slide.shapes:
            if shape.name == stat_shape_name or f"stat{i}" in shape.name.lower():
                if hasattr(shape, 'text_frame'):
                    # Combined format: "Label\nValue"
                    shape.text_frame.text = f"{label}\n{value}"
                    
                    # Apply formatting
                    if stat.get('color'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = hex_to_rgb(stat['color'])
                    
                    if stat.get('font_size'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.size = Pt(stat['font_size'])
                    
                    # Make value bold (second paragraph)
                    if len(shape.text_frame.paragraphs) > 1:
                        for run in shape.text_frame.paragraphs[1].runs:
                            run.font.bold = True
                    
                    print(f"📊 Updated Stat {i}: {label} = {value}")
                    found = True
                    break
        
        # If not found, try separate Label/Value shapes
        if not found:
            _update_separate_label_value(slide, i, label, value, stat)


def _update_separate_label_value(slide, index: int, label: str, value: str, stat: Dict[str, Any]):
    """Update separate label and value shapes"""
    label_shape_name = f"Label{index}"
    value_shape_name = f"Value{index}"
    
    # Update label
    for shape in slide.shapes:
        if shape.name == label_shape_name or f"label{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = label
                if stat.get('color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(stat['color'])
                break
    
    # Update value
    for shape in slide.shapes:
        if shape.name == value_shape_name or f"value{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = value
                if stat.get('color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(stat['color'])
                            run.font.bold = True
                
                if stat.get('font_size'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(stat['font_size'])
                
                print(f"📊 Updated Stat {index}: {label} = {value}")
                break


def _update_background_color(slide, background_color: str):
    """Update background color"""
    # This would need to be customized based on your template structure
    # For now, just log that we found background elements
    for shape in slide.shapes:
        if 'background' in shape.name.lower() or 'bg' in shape.name.lower():
            print(f"🎨 Found background element: {shape.name}")


def _update_image(slide, image_data: BytesIO):
    """Update image in slide"""
    for shape in slide.shapes:
        if 'image' in shape.name.lower() or 'picture' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(image_data, left, top, width, height)
                print(f"🖼️ Updated image")
                break
