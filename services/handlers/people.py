"""
People slide handler - Handles generation/modification of people/team slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_people_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle people/team slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - title: str - Slide title
            - description: str - Optional description
            - names: List[str] - List of people names
            - designations: List[str] - List of designations/titles
            - descriptions: List[str] - List of descriptions for each person
            - title_color: str - Color for title
            - description_color: str - Color for description
            - background_color: str - Color for background
        image_data: Optional BytesIO containing image data
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"👥 Processing PEOPLE slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update description
    if slide_data.get('description'):
        _update_description(slide, slide_data)
    
    # Update people data
    if slide_data.get('names'):
        _update_people(slide, slide_data)
    
    # Update background color if provided
    if slide_data.get('background_color'):
        _update_background_color(slide, slide_data['background_color'])
    
    # Update image if provided
    if image_data and slide_data.get('image_url'):
        _update_image(slide, image_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ People slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                if slide_data.get('title_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['title_color'])
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_description(slide, slide_data: Dict[str, Any]):
    """Update description text and formatting"""
    for shape in slide.shapes:
        if 'description' in shape.name.lower() or shape.name.startswith('Description'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['description']
                if slide_data.get('description_color'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['description_color'])
                print(f"📝 Updated description: {slide_data['description'][:50]}...")
                break


def _update_people(slide, slide_data: Dict[str, Any]):
    """Update people information"""
    names = slide_data.get('names', [])
    designations = slide_data.get('designations', [])
    descriptions = slide_data.get('descriptions', [])
    
    if not names:
        return
    
    # Ensure all arrays have the same length
    max_length = len(names)
    designations = designations + [''] * (max_length - len(designations))
    descriptions = descriptions + [''] * (max_length - len(descriptions))
    
    # Look for person shapes (Person1, Person2, etc.) or (Name1/Designation1/Description1, etc.)
    for i, (name, designation, description) in enumerate(zip(names, designations, descriptions), 1):
        # Try to find combined person shape first
        person_shape_name = f"Person{i}"
        found = False
        
        for shape in slide.shapes:
            if shape.name == person_shape_name or f"person{i}" in shape.name.lower():
                if hasattr(shape, 'text_frame'):
                    # Combined format: "Name\nDesignation\nDescription"
                    person_text = f"{name}\n{designation}"
                    if description:
                        person_text += f"\n{description}"
                    
                    shape.text_frame.text = person_text
                    
                    # Make name bold (first paragraph)
                    if shape.text_frame.paragraphs:
                        for run in shape.text_frame.paragraphs[0].runs:
                            run.font.bold = True
                    
                    # Make designation italic (second paragraph)
                    if len(shape.text_frame.paragraphs) > 1:
                        for run in shape.text_frame.paragraphs[1].runs:
                            run.font.italic = True
                    
                    print(f"👤 Updated Person {i}: {name} - {designation}")
                    found = True
                    break
        
        # If not found, try separate Name/Designation/Description shapes
        if not found:
            _update_separate_person_fields(slide, i, name, designation, description)


def _update_separate_person_fields(slide, index: int, name: str, designation: str, description: str):
    """Update separate name, designation, and description shapes"""
    name_shape_name = f"Name{index}"
    designation_shape_name = f"Designation{index}"
    description_shape_name = f"Description{index}"
    
    # Update name
    for shape in slide.shapes:
        if shape.name == name_shape_name or f"name{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = name
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
                break
    
    # Update designation
    for shape in slide.shapes:
        if shape.name == designation_shape_name or f"designation{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = designation
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.italic = True
                break
    
    # Update description
    for shape in slide.shapes:
        if shape.name == description_shape_name or f"description{index}" in shape.name.lower() or f"desc{index}" in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = description
                print(f"👤 Updated Person {index}: {name} - {designation}")
                break


def _update_background_color(slide, background_color: str):
    """Update background color"""
    # This would need to be customized based on your template structure
    # For now, just log that we found background elements
    for shape in slide.shapes:
        if 'background' in shape.name.lower() or 'bg' in shape.name.lower():
            print(f"🎨 Found background element: {shape.name}")


def _update_image(slide, image_data: BytesIO):
    """Update image in slide"""
    for shape in slide.shapes:
        if 'image' in shape.name.lower() or 'picture' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(image_data, left, top, width, height)
                print(f"🖼️ Updated image")
                break
