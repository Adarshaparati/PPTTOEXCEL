"""
Phases slide handler - Handles generation/modification of phase/timeline slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_phases_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle phases/timeline slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - title: str - Slide title
            - phases: List[Dict] - List of phases with name, description, status
            - timeline_color: str - Color for timeline elements
        image_data: Optional BytesIO containing image data
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"📅 Processing PHASES slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update phases
    if slide_data.get('phases'):
        _update_phases(slide, slide_data)
    
    # Update image if provided
    if image_data and slide_data.get('image_url'):
        _update_image(slide, image_data)
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Phases slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_phases(slide, slide_data: Dict[str, Any]):
    """Update phase information"""
    phases = slide_data.get('phases', [])
    if not phases:
        return
    
    # Look for phase text shapes (Phase1, Phase2, etc.)
    for i, phase in enumerate(phases, 1):
        phase_shape_name = f"Phase{i}"
        
        for shape in slide.shapes:
            if shape.name == phase_shape_name or f"phase{i}" in shape.name.lower():
                if hasattr(shape, 'text_frame'):
                    # Format: "Phase Name\nDescription\nStatus"
                    phase_text = f"{phase.get('name', '')}\n{phase.get('description', '')}"
                    if phase.get('status'):
                        phase_text += f"\nStatus: {phase['status']}"
                    
                    shape.text_frame.text = phase_text
                    
                    # Apply phase-specific formatting
                    if phase.get('color'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.font.color.rgb = hex_to_rgb(phase['color'])
                    
                    # Bold the phase name (first line)
                    if shape.text_frame.paragraphs:
                        first_paragraph = shape.text_frame.paragraphs[0]
                        for run in first_paragraph.runs:
                            run.font.bold = True
                    
                    print(f"📋 Updated Phase {i}: {phase.get('name', '')}")
                    break
    
    # Update timeline color if specified
    if slide_data.get('timeline_color'):
        _update_timeline_color(slide, slide_data['timeline_color'])


def _update_timeline_color(slide, timeline_color: str):
    """Update timeline visual elements color"""
    for shape in slide.shapes:
        if 'timeline' in shape.name.lower() or 'arrow' in shape.name.lower():
            # This would need to be customized based on your template structure
            # For now, just log that we found timeline elements
            print(f"🎨 Found timeline element: {shape.name}")


def _update_image(slide, image_data: BytesIO):
    """Update image in slide"""
    for shape in slide.shapes:
        if 'image' in shape.name.lower() or 'picture' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                sp = shape.element
                sp.getparent().remove(sp)
                slide.shapes.add_picture(image_data, left, top, width, height)
                print(f"🖼️ Updated image")
                break