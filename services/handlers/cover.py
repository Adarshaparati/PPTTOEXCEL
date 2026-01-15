"""
Cover slide handler - Handles generation/modification of cover slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any
import requests


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_cover_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle cover slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - title: str - Main title of the presentation
            - subtitle: str - Subtitle text
            - slide_name: str - Name of the slide
            - slide_data_id: str - Slide data identifier
            - slide_type: str - Type of slide (e.g., "Cover")
            - company_name: str - Company name
            - image: List[str] - Array of image URLs
            - colors: Dict[str, str] - Color scheme dictionary
        image_data: Optional BytesIO containing image data (not used directly, images come from URLs)
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"📄 Processing COVER slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update subtitle
    if slide_data.get('subtitle'):
        _update_subtitle(slide, slide_data)
    
    # Update company name
    if slide_data.get('company_name'):
        _update_company_name(slide, slide_data)
    
    # Update images from URLs
    if slide_data.get('image'):
        _update_images_from_urls(slide, slide_data.get('image', []))
    
    # Apply color scheme
    if slide_data.get('colors'):
        _apply_color_scheme(slide, slide_data.get('colors', {}))
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Cover slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                
                # Apply primary color if available
                if slide_data.get('colors', {}).get('primary'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['colors']['primary'])
                
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_subtitle(slide, slide_data: Dict[str, Any]):
    """Update subtitle text and formatting"""
    for shape in slide.shapes:
        if 'subtitle' in shape.name.lower() or shape.name.startswith('Subtitle'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['subtitle']
                
                # Apply secondary color if available
                if slide_data.get('colors', {}).get('secondary'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['colors']['secondary'])
                
                print(f"📝 Updated subtitle: {slide_data['subtitle']}")
                break


def _update_company_name(slide, slide_data: Dict[str, Any]):
    """Update company name text and formatting"""
    for shape in slide.shapes:
        if 'company' in shape.name.lower() or shape.name.startswith('Company'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['company_name']
                print(f"📝 Updated company name: {slide_data['company_name']}")
                break


def _update_images_from_urls(slide, image_urls: list):
    """Update images from URLs"""
    if not image_urls:
        return
    
    # Try to match images with shapes
    image_shapes = []
    for shape in slide.shapes:
        if 'image' in shape.name.lower() or 'logo' in shape.name.lower() or 'picture' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                image_shapes.append(shape)
    
    # Update images
    for i, (shape, image_url) in enumerate(zip(image_shapes, image_urls)):
        try:
            print(f"🖼️ Downloading image {i+1} from: {image_url}")
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            image_data = BytesIO(response.content)
            
            # Get position and size
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            # Remove old image
            sp = shape.element
            sp.getparent().remove(sp)
            
            # Add new image
            slide.shapes.add_picture(image_data, left, top, width, height)
            print(f"✅ Updated image {i+1}")
            
        except Exception as e:
            print(f"⚠️ Failed to update image {i+1} from {image_url}: {e}")


def _apply_color_scheme(slide, colors: Dict[str, str]):
    """Apply color scheme to slide elements"""
    # This would need to be customized based on your template structure
    # For now, just log that we have colors available
    if colors:
        print(f"🎨 Color scheme available: {', '.join(colors.keys())}")
        
        # Apply background color if specified
        if colors.get('background'):
            for shape in slide.shapes:
                if 'background' in shape.name.lower() or 'bg' in shape.name.lower():
                    print(f"🎨 Found background element: {shape.name}")
