"""
Contact slide handler - Handles generation/modification of contact slides
"""

from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from typing import Dict, Any
import requests


def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)
    
    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)
    
    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


def handle_contact_slide(presentation_bytes: BytesIO, slide_data: Dict[str, Any], image_data: BytesIO = None) -> BytesIO:
    """
    Handle contact slide modification
    
    Args:
        presentation_bytes: BytesIO containing the presentation
        slide_data: Dictionary with slide configuration
            - slide_number: int - Slide number to update
            - title: str - Contact slide title
            - slide_name: str - Name of the slide
            - website_link: str - Website URL
            - linkedin_link: str - LinkedIn profile URL
            - contact_email: str - Contact email address
            - contact_phone: str - Contact phone number
            - image: List[str] - Array of image URLs
            - colors: Dict[str, str] - Color scheme dictionary
        image_data: Optional BytesIO containing image data (not used directly, images come from URLs)
        
    Returns:
        BytesIO containing the modified presentation
    """
    print(f"📞 Processing CONTACT slide...")
    
    prs = Presentation(presentation_bytes)
    
    slide_index = slide_data.get('slide_number', 1) - 1
    if slide_index >= len(prs.slides):
        raise ValueError(f"Slide {slide_data.get('slide_number')} not found in presentation")
    
    slide = prs.slides[slide_index]
    
    # Update title
    if slide_data.get('title'):
        _update_title(slide, slide_data)
    
    # Update contact information
    if slide_data.get('website_link'):
        _update_website(slide, slide_data)
    
    if slide_data.get('linkedin_link'):
        _update_linkedin(slide, slide_data)
    
    if slide_data.get('contact_email'):
        _update_email(slide, slide_data)
    
    if slide_data.get('contact_phone'):
        _update_phone(slide, slide_data)
    
    # Update images from URLs
    if slide_data.get('image'):
        _update_images_from_urls(slide, slide_data.get('image', []))
    
    # Apply color scheme
    if slide_data.get('colors'):
        _apply_color_scheme(slide, slide_data.get('colors', {}))
    
    output = BytesIO()
    prs.save(output)
    output.seek(0)
    
    print(f"✅ Contact slide processed successfully")
    return output


def _update_title(slide, slide_data: Dict[str, Any]):
    """Update title text and formatting"""
    for shape in slide.shapes:
        if 'title' in shape.name.lower() or shape.name.startswith('Title'):
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['title']
                
                # Apply primary color if available
                if slide_data.get('colors', {}).get('primary'):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(slide_data['colors']['primary'])
                
                print(f"📝 Updated title: {slide_data['title']}")
                break


def _update_website(slide, slide_data: Dict[str, Any]):
    """Update website link"""
    for shape in slide.shapes:
        if 'website' in shape.name.lower() or 'web' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['website_link']
                
                # Add hyperlink if possible
                if shape.text_frame.paragraphs:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                run.hyperlink.address = slide_data['website_link']
                            except:
                                pass  # Hyperlink might not be supported
                
                print(f"🌐 Updated website: {slide_data['website_link']}")
                break


def _update_linkedin(slide, slide_data: Dict[str, Any]):
    """Update LinkedIn link"""
    for shape in slide.shapes:
        if 'linkedin' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['linkedin_link']
                
                # Add hyperlink if possible
                if shape.text_frame.paragraphs:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                run.hyperlink.address = slide_data['linkedin_link']
                            except:
                                pass
                
                print(f"🔗 Updated LinkedIn: {slide_data['linkedin_link']}")
                break


def _update_email(slide, slide_data: Dict[str, Any]):
    """Update email address"""
    for shape in slide.shapes:
        if 'email' in shape.name.lower() or 'mail' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['contact_email']
                
                # Add mailto link if possible
                if shape.text_frame.paragraphs:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            try:
                                run.hyperlink.address = f"mailto:{slide_data['contact_email']}"
                            except:
                                pass
                
                print(f"📧 Updated email: {slide_data['contact_email']}")
                break


def _update_phone(slide, slide_data: Dict[str, Any]):
    """Update phone number"""
    for shape in slide.shapes:
        if 'phone' in shape.name.lower() or 'tel' in shape.name.lower():
            if hasattr(shape, 'text_frame'):
                shape.text_frame.text = slide_data['contact_phone']
                print(f"📱 Updated phone: {slide_data['contact_phone']}")
                break


def _update_images_from_urls(slide, image_urls: list):
    """Update images from URLs"""
    if not image_urls:
        return
    
    # Try to match images with shapes
    image_shapes = []
    for shape in slide.shapes:
        if 'image' in shape.name.lower() or 'qr' in shape.name.lower() or 'picture' in shape.name.lower():
            if shape.shape_type == 13:  # Picture type
                image_shapes.append(shape)
    
    # Update images
    for i, (shape, image_url) in enumerate(zip(image_shapes, image_urls)):
        try:
            print(f"🖼️ Downloading image {i+1} from: {image_url}")
            response = requests.get(image_url, timeout=10)
            response.raise_for_status()
            
            image_data = BytesIO(response.content)
            
            # Get position and size
            left, top, width, height = shape.left, shape.top, shape.width, shape.height
            
            # Remove old image
            sp = shape.element
            sp.getparent().remove(sp)
            
            # Add new image
            slide.shapes.add_picture(image_data, left, top, width, height)
            print(f"✅ Updated image {i+1}")
            
        except Exception as e:
            print(f"⚠️ Failed to update image {i+1} from {image_url}: {e}")


def _apply_color_scheme(slide, colors: Dict[str, str]):
    """Apply color scheme to slide elements"""
    # This would need to be customized based on your template structure
    # For now, just log that we have colors available
    if colors:
        print(f"🎨 Color scheme available: {', '.join(colors.keys())}")
        
        # Apply text color if specified
        if colors.get('text'):
            text_color = hex_to_rgb(colors['text'])
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame') and 'title' not in shape.name.lower():
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = text_color
