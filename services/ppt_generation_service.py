import os
import io
import requests
from io import BytesIO
from pptx import Presentation
from pptx.dml.color import RGBColor
from typing import Dict, List, Any
from services.sheets_service import get_sheets_service
from services.s3_service import s3_service
from datetime import datetime

class PPTGenerationService:
    """Service for generating PowerPoint presentations from templates"""
    
    @staticmethod
    def hex_to_rgb(hex_str):
        """Convert #RRGGBB string to RGBColor safely."""
        if not hex_str:
            return RGBColor(0, 0, 0)
        
        hex_str = str(hex_str).strip().replace("#", "")
        if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
            print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
            return RGBColor(0, 0, 0)
        
        return RGBColor(int(hex_str[0:2], 16),
                        int(hex_str[2:4], 16),
                        int(hex_str[4:6], 16))
    
    def get_template_data_from_sheets(self, template_id: str) -> List[Dict]:
        """
        Fetch all rows from Google Sheets matching the template_id
        
        Args:
            template_id: Template identifier to filter data
            
        Returns:
            List of dictionaries with shape data
        """
        sheets_svc = get_sheets_service()
        if not sheets_svc:
            raise RuntimeError("Could not initialize Google Sheets service")
        
        # Get all data from sheet
        result = sheets_svc.get_sheet_data("Sheet1")
        if not result.get('success'):
            raise RuntimeError(f"Failed to fetch sheet data: {result.get('error')}")
        
        rows = result.get('data', [])
        if not rows or len(rows) < 2:
            raise RuntimeError("No data found in Google Sheets")
        
        # First row is headers
        headers = rows[0]
        data_rows = rows[1:]
        
        # Filter rows by template_id (column A)
        filtered_data = []
        for row in data_rows:
            if len(row) > 0 and row[0] == template_id:
                # Convert row to dictionary using headers
                row_dict = {}
                for i, header in enumerate(headers):
                    row_dict[header] = row[i] if i < len(row) else ""
                filtered_data.append(row_dict)
        
        print(f"📊 Found {len(filtered_data)} rows for template_id: {template_id}")
        return filtered_data
    
    def build_replacements_dict(self, sheet_data: List[Dict]) -> Dict[str, str]:
        """
        Build a dictionary of shape_name → content mappings
        Maps actual shape names from PPT to their content values
        
        Args:
            sheet_data: List of row dictionaries from Google Sheets
            
        Returns:
            Dictionary mapping shape names to content values
        """
        replacements = {}
        
        # Map each shape name to its sub-component content
        for row in sheet_data:
            shape_name = row.get('Shape Name', '').strip()
            
            # Get sub-component value from "Component | Sub-component" column
            sub_component = row.get('Sub-component', '').strip()
            
            # Fallback to Content column if Sub-component is empty
            if not sub_component:
                sub_component = row.get('Content', '').strip()
            
            if not shape_name:
                continue
            
            # Direct mapping: shape name → sub-component value
            replacements[shape_name] = sub_component
            
            # Also create lowercase version for case-insensitive matching
            replacements[shape_name.lower()] = sub_component
            
            print(f"  📝 '{shape_name}' → '{sub_component}'")
        
        # Add color placeholders if they exist
        if sheet_data:
            # Get unique colors from all rows
            fill_colors = set()
            font_colors = set()
            
            for row in sheet_data:
                fill_color = row.get('Fill Color', '').strip()
                font_color = row.get('Font Color', '').strip()
                
                if fill_color and fill_color.startswith('#'):
                    fill_colors.add(fill_color)
                if font_color and font_color.startswith('#'):
                    font_colors.add(font_color)
            
            # Use first found colors as primary/secondary
            if fill_colors:
                replacements["{{P100}}"] = list(fill_colors)[0]
            if font_colors:
                replacements["{{S100}}"] = list(font_colors)[0]
        
        print(f"🔑 Built {len(replacements)} shape name → sub-component mappings")
        # Show first 5 mappings as examples
        for i, (name, value) in enumerate(list(replacements.items())[:5]):
            print(f"   '{name}' → '{value}'")
        return replacements
    
    def replace_text_in_presentation(self, prs: Presentation, replacements: Dict[str, str]) -> int:
        """
        Replace shape text based on shape names matching the data
        
        Args:
            prs: PowerPoint presentation object
            replacements: Dictionary of shape_name → content mappings
            
        Returns:
            Number of replacements made
        """
        replacement_count = 0
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                # Get the shape name
                shape_name = shape.name.strip()
                shape_name_lower = shape_name.lower()
                
                # Check if this shape name exists in our mappings
                content = None
                if shape_name in replacements:
                    content = replacements[shape_name]
                elif shape_name_lower in replacements:
                    content = replacements[shape_name_lower]
                
                # Replace the text if we found content
                if content is not None:
                    if shape.has_text_frame:
                        # Replace all text in the shape with content
                        original_text = shape.text
                        shape.text = str(content)
                        replacement_count += 1
                        print(f"  ✓ Slide {slide_idx}: Shape '{shape_name}' → '{content[:50]}...'")
                    elif hasattr(shape, 'text'):
                        try:
                            shape.text = str(content)
                            replacement_count += 1
                            print(f"  ✓ Slide {slide_idx}: Shape '{shape_name}' → '{content[:50]}...'")
                        except:
                            pass
                
                # Also check for {{placeholder}} patterns in existing text
                if shape.has_text_frame:
                    original_text = shape.text
                    new_text = original_text
                    
                    # Replace any {{ShapeName}} placeholders
                    for shape_key, value in replacements.items():
                        placeholder = f"{{{{{shape_key}}}}}"
                        if placeholder in new_text:
                            new_text = new_text.replace(placeholder, str(value))
                            replacement_count += 1
                            print(f"  ✓ Slide {slide_idx}: Replaced placeholder {placeholder[:30]}... → '{value[:50]}...'")
                    
                    # Update if changed
                    if new_text != original_text:
                        shape.text = new_text
        
        print(f"✅ Made {replacement_count} text replacements")
        return replacement_count
    
    def apply_colors_from_data(self, prs: Presentation, sheet_data: List[Dict]):
        """
        Apply fill colors and font colors based on sheet data
        
        Args:
            prs: PowerPoint presentation object
            sheet_data: List of row dictionaries with color information
        """
        # Build color map by shape name
        color_map = {}
        for row in sheet_data:
            shape_name = row.get('Shape Name', '').lower().strip()
            fill_color = row.get('Fill Color', '')
            font_color = row.get('Font Color', '')
            
            if shape_name:
                color_map[shape_name] = {
                    'fill': fill_color,
                    'font': font_color
                }
        
        # Apply colors to shapes
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                shape_key = shape.name.lower().strip()
                
                if shape_key in color_map:
                    colors = color_map[shape_key]
                    
                    # Apply fill color
                    if colors['fill']:
                        try:
                            if hasattr(shape, 'fill'):
                                shape.fill.solid()
                                shape.fill.fore_color.rgb = self.hex_to_rgb(colors['fill'])
                                print(f"  🎨 Applied fill color to '{shape.name}'")
                        except Exception as e:
                            print(f"  ⚠️ Could not apply fill color to '{shape.name}': {e}")
                    
                    # Apply font color
                    if colors['font'] and shape.has_text_frame:
                        try:
                            rgb_color = self.hex_to_rgb(colors['font'])
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = rgb_color
                            print(f"  🎨 Applied font color to '{shape.name}'")
                        except Exception as e:
                            print(f"  ⚠️ Could not apply font color to '{shape.name}': {e}")
    
    def replace_images_from_data(self, prs: Presentation, sheet_data: List[Dict]):
        """
        Replace images in shapes based on Image S3 URL from sheet data
        
        Args:
            prs: PowerPoint presentation object
            sheet_data: List of row dictionaries with image URLs
        """
        # Build image map by shape name
        image_map = {}
        for row in sheet_data:
            shape_name = row.get('Shape Name', '').lower().strip()
            image_url = row.get('Image S3 URL', '') or row.get('Image URL', '')
            
            if shape_name and image_url and image_url.startswith('http'):
                image_map[shape_name] = image_url
        
        # Replace images
        for slide_idx, slide in enumerate(prs.slides, start=1):
            for shape in slide.shapes:
                shape_key = shape.name.lower().strip()
                
                if shape_key in image_map:
                    image_url = image_map[shape_key]
                    try:
                        # Download image
                        print(f"  📥 Downloading image for '{shape.name}' from {image_url[:50]}...")
                        response = requests.get(image_url, timeout=10)
                        response.raise_for_status()
                        image_bytes = BytesIO(response.content)
                        
                        # Replace shape with image
                        left, top, width, height = shape.left, shape.top, shape.width, shape.height
                        slide.shapes._spTree.remove(shape._element)
                        slide.shapes.add_picture(image_bytes, left, top, width, height)
                        print(f"  🖼️ Image replaced in '{shape.name}'")
                    except Exception as e:
                        print(f"  ⚠️ Could not replace image in '{shape.name}': {e}")
    
    def generate_presentation(
        self,
        template_s3_url: str,
        template_id: str,
        output_filename: str = None
    ) -> tuple:
        """
        Generate a PowerPoint presentation from template and Google Sheets data
        
        Args:
            template_s3_url: S3 URL of the template PowerPoint file
            template_id: Template ID to fetch data from Google Sheets
            output_filename: Optional custom filename for output
            
        Returns:
            tuple: (presentation_bytes, filename, stats_dict)
        """
        print(f"\n🚀 Starting presentation generation...")
        print(f"📄 Template: {template_s3_url}")
        print(f"🆔 Template ID: {template_id}")
        
        # Download template from S3
        from urllib.parse import urlparse, unquote
        parsed_url = urlparse(template_s3_url)
        s3_key = unquote(parsed_url.path.lstrip('/'))
        
        print(f"📥 Downloading template from S3: {s3_key}")
        template_bytes = s3_service.download_file(s3_key)
        if not template_bytes:
            raise RuntimeError(f"Failed to download template from S3: {s3_key}")
        
        # Load presentation
        prs = Presentation(BytesIO(template_bytes))
        print(f"✅ Template loaded: {len(prs.slides)} slides")
        
        # Get data from Google Sheets
        sheet_data = self.get_template_data_from_sheets(template_id)
        if not sheet_data:
            raise RuntimeError(f"No data found for template_id: {template_id}")
        
        # Build replacement mappings
        replacements = self.build_replacements_dict(sheet_data)
        
        # Apply all transformations
        print("\n🔄 Applying transformations...")
        text_replacements = self.replace_text_in_presentation(prs, replacements)
        
        print("\n🎨 Applying colors...")
        self.apply_colors_from_data(prs, sheet_data)
        
        print("\n🖼️ Replacing images...")
        self.replace_images_from_data(prs, sheet_data)
        
        # Generate output filename
        if not output_filename:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_filename = f"Generated_{template_id}_{timestamp}.pptx"
        
        # Save to bytes
        output_buffer = BytesIO()
        prs.save(output_buffer)
        output_buffer.seek(0)
        presentation_bytes = output_buffer.getvalue()
        
        stats = {
            'template_id': template_id,
            'slides_count': len(prs.slides),
            'data_rows': len(sheet_data),
            'text_replacements': text_replacements,
            'output_filename': output_filename,
            'file_size': len(presentation_bytes)
        }
        
        print(f"\n✅ Presentation generated successfully!")
        print(f"   📊 {stats['slides_count']} slides")
        print(f"   📝 {stats['text_replacements']} text replacements")
        print(f"   💾 {stats['file_size'] / 1024:.2f} KB")
        
        return presentation_bytes, output_filename, stats


# Create singleton instance
ppt_generation_service = PPTGenerationService()
