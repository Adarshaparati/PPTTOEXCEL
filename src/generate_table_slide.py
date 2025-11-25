import os
import requests
from io import BytesIO
from msal import PublicClientApplication
from dotenv import load_dotenv
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor
from datetime import datetime

# ---------------------
# Load environment
# ---------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "").strip()
CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
AUTHORITY = os.getenv("AUTHORITY", "").strip() or f"https://login.microsoftonline.com/{TENANT_ID}"

EXCEL_ONEDRIVE_PATH = os.getenv("EXCEL_ONEDRIVE_PATH", "/me/drive/root:/Book.xlsx")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"


# ---------------------
# Authentication
# ---------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]


# ---------------------
# Graph helpers
# ---------------------
def graph_get(url, token, stream=False):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers, stream=stream)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r


def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r


# ---------------------
# Read Excel (Sheet3)
# ---------------------
def read_excel_values(token):
    file_url = f"{GRAPH_ROOT}{EXCEL_ONEDRIVE_PATH}:/content"
    print(f"📥 Downloading Excel from {EXCEL_ONEDRIVE_PATH} ...")
    r = graph_get(file_url, token, stream=True)

    wb = load_workbook(filename=BytesIO(r.content), data_only=True)
    sheet = wb["Sheet3"]

    data = {
        "Slide_No": int(sheet["A2"].value or 12),
        "Title": sheet["B2"].value or "Business Overview",
        "RowHeader1": sheet["C2"].value or "Services",
        "RowHeader2": sheet["D2"].value or "Operations",
        "Column_Header1": sheet["E2"].value or "Revenue",
        "Column_Header2": sheet["F2"].value or "Growth",
        "P100": sheet["G2"].value or "#3667B2",
        "S100": sheet["H2"].value or "#8A8B8C",
        "C1R1": sheet["I2"].value or "",
        "C2R1": sheet["J2"].value or "",
        "C1R2": sheet["K2"].value or "",
        "C2R2": sheet["L2"].value or "",
    }

    wb.close()
    print("📘 Excel Data Loaded:")
    for k, v in data.items():
        print(f"   {k}: {v}")
    return data


# ---------------------
# Update PPT
# ---------------------
def update_table_slide(template_bytes, values):
    prs = Presentation(BytesIO(template_bytes))
    slide_index = values["Slide_No"] - 1
    if slide_index >= len(prs.slides):
        raise IndexError(f"Slide {values['Slide_No']} not found.")
    slide = prs.slides[slide_index]

    def hex_to_rgb(hex_str):
        hex_str = hex_str.lstrip("#")
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    p100_color = hex_to_rgb(values["P100"])
    s100_color = hex_to_rgb(values["S100"])

    # Replace text placeholders and apply colors
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text
            text = text.replace("{{Title}}", str(values["Title"]))
            text = text.replace("{{RowHeader1}}", str(values["RowHeader1"]))
            text = text.replace("{{RowHeader2}}", str(values["RowHeader2"]))
            text = text.replace("{{RowHeader}}", str(values["RowHeader1"]))
            text = text.replace("{{Column_Header1}}", str(values["Column_Header1"]))
            text = text.replace("{{Column_Header2}}", str(values["Column_Header2"]))
            text = text.replace("{{Column_Header}}", str(values["Column_Header1"]))

            # Replace value placeholders
            text = text.replace("{{C1R1_VALUE}}", str(values["C1R1"]))
            text = text.replace("{{C2R1_VALUE}}", str(values["C2R1"]))
            text = text.replace("{{C1R2_VALUE}}", str(values["C1R2"]))
            text = text.replace("{{C2R2_VALUE}}", str(values["C2R2"]))
            text = text.replace("{{value1}}", "")  # clean leftover placeholders
            shape.text = text

            # Set text color
            for p in shape.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = s100_color

        # Column Header fill → P100
        elif "Column_Header" in shape.name:
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = p100_color
            except:
                pass

        # Row Header fill → S100
        elif "RowHeader" in shape.name:
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = s100_color
            except:
                pass

        # Image backgrounds → S100
        elif "image_bg_1" in shape.name.lower() or "image_bg_2" in shape.name.lower():
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = s100_color
                print(f"🖼️ {shape.name} recolored to S100 {values['S100']}")
            except:
                pass

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    print("✅ Table slide updated successfully.")
    return bio


# ---------------------
# Main
# ---------------------
def main():
    token = acquire_token_device_code()
    values = read_excel_values(token)
    r = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token, stream=True)
    template_bytes = r.content

    updated_ppt = update_table_slide(template_bytes, values)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_title = "".join(c for c in values["Title"] if c.isalnum() or c in (" ", "_", "-")).strip()
    dest_path = f"{DEST_FOLDER_ONEDRIVE}/{safe_title.replace(' ', '_')}_Table_{timestamp}.pptx"
    upload_url = f"{GRAPH_ROOT}{dest_path}:/content?@microsoft.graph.conflictBehavior=rename"

    print(f"📤 Uploading to {dest_path} ...")
    graph_put(upload_url, token, data=updated_ppt.getvalue())
    print("✅ Upload complete!")


if __name__ == "__main__":
    main()
