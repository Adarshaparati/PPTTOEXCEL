import os
import requests
from io import BytesIO
from datetime import datetime
from dotenv import load_dotenv
from msal import PublicClientApplication
from pptx import Presentation
from openpyxl import Workbook
import base64
import hashlib
try:
    from PIL import Image
except ImportError:
    Image = None  # PIL is optional for enhanced image analysis

# ------------------------------------------------
# Load environment
# ------------------------------------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "")
CLIENT_ID = os.getenv("CLIENT_ID", "")
AUTHORITY = os.getenv("AUTHORITY", f"https://login.microsoftonline.com/{TENANT_ID}")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# ------------------------------------------------
# Auth helper
# ------------------------------------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]

# ------------------------------------------------
# Graph helpers
# ------------------------------------------------
def graph_get(url, token):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r

def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r

def get_image_extension(content_type):
    """Get file extension from content type"""
    extension_map = {
        'image/jpeg': '.jpg',
        'image/jpg': '.jpg', 
        'image/png': '.png',
        'image/gif': '.gif',
        'image/bmp': '.bmp',
        'image/tiff': '.tiff',
        'image/webp': '.webp',
        'image/svg+xml': '.svg'
    }
    return extension_map.get(content_type.lower(), '.img')

def get_mime_type(content_type):
    """Get proper MIME type for base64 data URL"""
    if content_type and content_type.startswith('image/'):
        return content_type
    return 'image/jpeg'  # Default fallback

def sanitize_text(text):
    """Remove illegal characters for Excel cells"""
    if text is None:
        return ''
    if not isinstance(text, str):
        text = str(text)
    # Remove control characters except tab, newline, and carriage return
    # Excel doesn't support characters with codes 0-31 except 9 (tab), 10 (LF), 13 (CR)
    import re
    # Remove all control characters except allowed ones
    text = re.sub(r'[\x00-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', text)
    # Limit length to avoid Excel's 32,767 character limit per cell
    if len(text) > 32000:
        text = text[:32000] + '...'
    return text

# ------------------------------------------------
# Extract presentation data
# ------------------------------------------------
def extract_ppt_to_excel(prs):
    wb = Workbook()
    ws = wb.active
    ws.title = "PPT_Data"
    
    # Enhanced headers for comprehensive data
    headers = [
        "Slide No", "Shape Name", "Shape Type", "Content", 
        "Left (EMU)", "Top (EMU)", "Width (EMU)", "Height (EMU)",
        "Left (Inches)", "Top (Inches)", "Width (Inches)", "Height (Inches)",
        "Font Name", "Font Size", "Font Bold", "Font Italic", "Font Underline", "Font Color",
        "Text Alignment", "Line Spacing", "Paragraph Spacing",
        "Fill Color", "Fill Type", "Transparency", "Line Color", "Line Width", "Line Style", "Rotation",
        "Has Image", "Image Format", "Image Width", "Image Height", "Image File Size", "Image URL", "Image Base64",
        "Chart Type", "Chart Title", "Chart Data", "Chart Categories", "Chart Series",
        "Hyperlink", "Z-Order", "Hidden", "Shadow", "Glow Effect", "Reflection",
        "3D Effects", "Placeholder Type", "Animation Effects"
    ]
    ws.append(headers)

    def emu_to_inches(emu):
        """Convert EMU (English Metric Units) to inches"""
        return round(emu / 914400, 2) if emu else 0

    def get_color_info(color_obj):
        """Extract color information"""
        try:
            if hasattr(color_obj, 'rgb'):
                return f"RGB({color_obj.rgb.red},{color_obj.rgb.green},{color_obj.rgb.blue})"
            elif hasattr(color_obj, 'theme_color'):
                return f"Theme Color {color_obj.theme_color}"
            else:
                return "Auto/Default"
        except:
            return "Unknown"

    def get_font_info(text_frame):
        """Extract comprehensive font information from text frame"""
        font_info = {
            'name': 'Default',
            'size': 'Default',
            'bold': False,
            'italic': False,
            'underline': False,
            'color': 'Auto',
            'alignment': 'Left',
            'line_spacing': 'Default',
            'paragraph_spacing': 'Default'
        }
        
        try:
            if text_frame and text_frame.paragraphs:
                first_paragraph = text_frame.paragraphs[0]
                
                # Paragraph alignment
                if hasattr(first_paragraph, 'alignment'):
                    alignment_map = {
                        0: 'Left', 1: 'Center', 2: 'Right', 3: 'Justify', 
                        4: 'Distribute', 5: 'Thai Distribute'
                    }
                    font_info['alignment'] = alignment_map.get(first_paragraph.alignment, 'Left')
                
                # Line spacing
                if hasattr(first_paragraph, 'line_spacing'):
                    try:
                        if first_paragraph.line_spacing:
                            font_info['line_spacing'] = f"{first_paragraph.line_spacing:.2f}"
                    except:
                        pass
                
                # Space before/after paragraph
                try:
                    space_before = first_paragraph.space_before.pt if first_paragraph.space_before else 0
                    space_after = first_paragraph.space_after.pt if first_paragraph.space_after else 0
                    font_info['paragraph_spacing'] = f"Before:{space_before}pt, After:{space_after}pt"
                except:
                    pass
                
                if first_paragraph.runs:
                    first_run = first_paragraph.runs[0]
                    font = first_run.font
                    
                    font_info['name'] = font.name or 'Default'
                    font_info['size'] = f"{font.size.pt}pt" if font.size else 'Default'
                    font_info['bold'] = font.bold if font.bold is not None else False
                    font_info['italic'] = font.italic if font.italic is not None else False
                    font_info['underline'] = font.underline if font.underline is not None else False
                    font_info['color'] = get_color_info(font.color) if hasattr(font, 'color') else 'Auto'
        except:
            pass
        
        return font_info

    def get_fill_info(shape):
        """Extract detailed fill information"""
        fill_info = {
            'color': 'None',
            'type': 'None',
            'transparency': 0
        }
        
        try:
            if hasattr(shape, 'fill') and shape.fill:
                # Fill type
                fill_type_map = {
                    0: 'No Fill', 1: 'Solid', 2: 'Gradient', 3: 'Picture', 
                    4: 'Pattern', 5: 'Group', 6: 'Background'
                }
                fill_info['type'] = fill_type_map.get(shape.fill.type, 'Unknown')
                
                # Fill color
                if hasattr(shape.fill, 'fore_color'):
                    fill_info['color'] = get_color_info(shape.fill.fore_color)
                
                # Transparency
                if hasattr(shape.fill, 'transparency'):
                    fill_info['transparency'] = f"{shape.fill.transparency * 100:.1f}%"
                    
        except:
            pass
        
        return fill_info

    def get_image_info(shape, slide_num, shape_index, images_folder="extracted_images"):
        """Extract detailed image information and save images"""
        image_info = {
            'has_image': False,
            'format': '',
            'width': 0,
            'height': 0,
            'file_size': 0,
            'url': '',
            'base64': ''
        }
        
        try:
            if shape.shape_type == 12 or hasattr(shape, 'image'):  # Picture
                image_info['has_image'] = True
                
                if hasattr(shape, 'image') and shape.image:
                    # Get image data
                    image_blob = shape.image.blob
                    image_info['format'] = shape.image.content_type or 'Unknown'
                    image_info['file_size'] = len(image_blob) if image_blob else 0
                    
                    if image_blob:
                        # Create images folder if it doesn't exist
                        if not os.path.exists(images_folder):
                            os.makedirs(images_folder)
                        
                        # Generate unique filename based on content hash
                        image_hash = hashlib.md5(image_blob).hexdigest()[:12]
                        file_extension = get_image_extension(image_info['format'])
                        filename = f"slide_{slide_num}_shape_{shape_index}_{image_hash}{file_extension}"
                        filepath = os.path.join(images_folder, filename)
                        
                        # Save image to file
                        with open(filepath, 'wb') as f:
                            f.write(image_blob)
                        
                        # Create URL (relative path)
                        image_info['url'] = filepath.replace('\\', '/')
                        
                        # Create base64 data URL for direct embedding
                        base64_string = base64.b64encode(image_blob).decode('utf-8')
                        mime_type = get_mime_type(image_info['format'])
                        image_info['base64'] = f"data:{mime_type};base64,{base64_string}"  # Full base64 for complete usage
                    
                    # Try to get image dimensions
                    try:
                        if Image and image_blob:  # Only if PIL is available
                            img = Image.open(BytesIO(image_blob))
                            image_info['width'] = img.width
                            image_info['height'] = img.height
                        else:
                            # Fallback to shape dimensions
                            image_info['width'] = emu_to_inches(shape.width)
                            image_info['height'] = emu_to_inches(shape.height)
                    except:
                        # Fallback to shape dimensions
                        image_info['width'] = emu_to_inches(shape.width)
                        image_info['height'] = emu_to_inches(shape.height)
                        
        except:
            pass
        
        return image_info

    def get_chart_info(shape):
        """Extract detailed chart information"""
        chart_info = {
            'type': 'None',
            'title': 'None',
            'data': 'None',
            'categories': 'None',
            'series': 'None'
        }
        
        try:
            if shape.shape_type == 3 and hasattr(shape, 'chart'):  # Chart type
                chart = shape.chart
                
                # Chart type
                chart_type_map = {
                    1: 'Area', 2: 'Bar', 3: 'Column', 4: 'Line', 5: 'Pie',
                    6: 'Scatter', 7: 'Surface', 8: 'Radar', 9: 'Treemap',
                    10: 'Sunburst', 11: 'Histogram', 12: 'BoxWhisker',
                    13: 'Waterfall', 14: 'Funnel', 15: 'Map'
                }
                chart_info['type'] = chart_type_map.get(chart.chart_type, f'Unknown({chart.chart_type})')
                
                # Chart title
                if hasattr(chart, 'chart_title') and chart.chart_title:
                    try:
                        if hasattr(chart.chart_title, 'text_frame') and chart.chart_title.text_frame:
                            chart_info['title'] = chart.chart_title.text_frame.text.strip()
                    except:
                        chart_info['title'] = 'Has Title'
                
                # Extract series data
                series_data = []
                categories_data = []
                chart_data_points = []
                
                try:
                    if hasattr(chart, 'plots') and chart.plots:
                        plot = chart.plots[0]  # Get first plot
                        
                        # Extract categories
                        if hasattr(plot, 'categories') and plot.categories:
                            try:
                                categories_data = [str(cat) for cat in plot.categories]
                            except:
                                categories_data = ['Category data available']
                        
                        # Extract series
                        if hasattr(plot, 'series'):
                            for i, series in enumerate(plot.series):
                                series_name = getattr(series, 'name', f'Series {i+1}')
                                series_data.append(series_name)
                                
                                # Try to extract values
                                try:
                                    if hasattr(series, 'values') and series.values:
                                        values = [str(v) for v in series.values if v is not None]
                                        if values:
                                            chart_data_points.append(f"{series_name}: [{', '.join(values[:10])}{'...' if len(values) > 10 else ''}]")
                                except:
                                    chart_data_points.append(f"{series_name}: [Values available]")
                except:
                    pass
                
                # Format extracted data
                if categories_data:
                    chart_info['categories'] = ' | '.join(categories_data[:10]) + ('...' if len(categories_data) > 10 else '')
                
                if series_data:
                    chart_info['series'] = ' | '.join(series_data)
                
                if chart_data_points:
                    chart_info['data'] = ' | '.join(chart_data_points)
                
                # Fallback: Try to extract from chart XML
                if chart_info['data'] == 'None':
                    try:
                        chart_xml = shape.element.xml
                        if 'val' in chart_xml and 'cat' in chart_xml:
                            chart_info['data'] = '[Chart data embedded in XML]'
                    except:
                        pass
                        
        except Exception as e:
            # If chart extraction fails, at least indicate it's a chart
            if shape.shape_type == 3:
                chart_info['type'] = 'Chart (extraction failed)'
                chart_info['data'] = f'Chart present but data extraction failed: {str(e)[:50]}'
        
        return chart_info

    def get_effects_info(shape):
        """Extract visual effects information"""
        effects_info = {
            'shadow': False,
            'glow': False,
            'reflection': False,
            '3d_effects': False
        }
        
        try:
            # Check for shadow
            if hasattr(shape, 'shadow') and shape.shadow.inherit:
                effects_info['shadow'] = True
                
            # Check for other effects (these might not be directly accessible)
            # This is a simplified check - actual implementation might vary
            if hasattr(shape, 'element'):
                element_xml = shape.element.xml
                if 'glow' in element_xml.lower():
                    effects_info['glow'] = True
                if 'reflection' in element_xml.lower():
                    effects_info['reflection'] = True
                if 'scene3d' in element_xml.lower() or 'sp3d' in element_xml.lower():
                    effects_info['3d_effects'] = True
                    
        except:
            pass
        
        return effects_info

    def get_placeholder_type(shape):
        """Get placeholder type if shape is a placeholder"""
        try:
            if hasattr(shape, 'placeholder') and shape.placeholder:
                placeholder_map = {
                    0: 'Title', 1: 'Body', 2: 'CenterTitle', 3: 'Subtitle',
                    4: 'DateAndTime', 5: 'SlideNumber', 6: 'Footer', 7: 'Header',
                    8: 'Object', 9: 'Chart', 10: 'Table', 11: 'ClipArt',
                    12: 'Diagram', 13: 'Media', 14: 'SlideImage', 15: 'Picture'
                }
                return placeholder_map.get(shape.placeholder.placeholder_format.type, 'Unknown')
        except:
            pass
        return 'Not a placeholder'

    def get_shape_type_name(shape_type):
        """Convert shape type number to readable name"""
        type_mapping = {
            1: "AutoShape", 2: "Callout", 3: "Chart", 4: "Comment", 
            5: "Freeform", 6: "Group", 7: "Line", 8: "LinkedOLEObject",
            9: "LinkedPicture", 10: "Media", 11: "OLEObject", 12: "Picture",
            13: "Placeholder", 14: "TextBox", 15: "3DModel", 16: "Canvas",
            17: "Connector", 18: "Ink", 19: "Table", 20: "SmartArt"
        }
        return type_mapping.get(shape_type, f"Unknown({shape_type})")

    for slide_num, slide in enumerate(prs.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes):
            # Basic shape info
            content = ""
            hyperlink = ""
            
            # Extract content based on shape type
            if shape.has_text_frame:
                content = sanitize_text(shape.text.strip().replace("\n", " | "))
                # Check for hyperlinks
                try:
                    if shape.text_frame.paragraphs:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if hasattr(run, 'hyperlink') and run.hyperlink.address:
                                    hyperlink = sanitize_text(run.hyperlink.address)
                                    break
                except:
                    pass
                    
            elif shape.shape_type == 19:  # Table
                table_data = []
                try:
                    for row in shape.table.rows:
                        row_text = [sanitize_text(cell.text.strip()) for cell in row.cells]
                        table_data.append(", ".join(row_text))
                    content = " | ".join(table_data)
                except:
                    content = "[TABLE - Could not extract data]"
                    
            elif shape.shape_type == 12:  # Picture
                content = "[IMAGE]"
                
            elif shape.shape_type == 3:  # Chart
                content = "[CHART]"

            # Position and size info
            left_emu = shape.left
            top_emu = shape.top
            width_emu = shape.width
            height_emu = shape.height
            
            left_inches = emu_to_inches(left_emu)
            top_inches = emu_to_inches(top_emu)
            width_inches = emu_to_inches(width_emu)
            height_inches = emu_to_inches(height_emu)

            # Enhanced information extraction
            font_info = get_font_info(shape.text_frame if shape.has_text_frame else None)
            fill_info = get_fill_info(shape)
            image_info = get_image_info(shape, slide_num, shape_index)
            chart_info = get_chart_info(shape)
            effects_info = get_effects_info(shape)
            placeholder_type = get_placeholder_type(shape)

            # Line information
            line_color = "None"
            line_width = "None"
            line_style = "None"
            
            try:
                if hasattr(shape, 'line') and shape.line:
                    line_color = get_color_info(shape.line.color)
                    line_width = f"{shape.line.width.pt}pt" if shape.line.width else "Default"
                    # Line style
                    line_style_map = {
                        0: 'None', 1: 'Solid', 2: 'Dash', 3: 'Dot', 
                        4: 'DashDot', 5: 'DashDotDot', 6: 'Double'
                    }
                    if hasattr(shape.line, 'dash_style'):
                        line_style = line_style_map.get(shape.line.dash_style, 'Solid')
            except:
                pass

            # Rotation
            rotation = 0
            try:
                rotation = shape.rotation
            except:
                pass

            # Shape type
            shape_type_name = get_shape_type_name(shape.shape_type)
            
            # Hidden status
            hidden = False
            try:
                hidden = not shape.element.get('hidden', '0') == '0'
            except:
                pass

            # Animation effects (simplified check)
            animation_effects = "None"
            try:
                # This is a basic check - full animation detection would require more complex parsing
                if hasattr(shape, 'element') and 'anim' in shape.element.xml.lower():
                    animation_effects = "Has Animation"
            except:
                pass

            # Append all enhanced data to worksheet (sanitize all string values)
            row_data = [
                slide_num, sanitize_text(shape.name), shape_type_name, content,
                left_emu, top_emu, width_emu, height_emu,
                left_inches, top_inches, width_inches, height_inches,
                sanitize_text(font_info['name']), sanitize_text(font_info['size']), font_info['bold'], 
                font_info['italic'], font_info['underline'], sanitize_text(font_info['color']),
                sanitize_text(font_info['alignment']), sanitize_text(font_info['line_spacing']), sanitize_text(font_info['paragraph_spacing']),
                sanitize_text(fill_info['color']), sanitize_text(fill_info['type']), sanitize_text(fill_info['transparency']),
                sanitize_text(line_color), sanitize_text(line_width), sanitize_text(line_style), rotation,
                image_info['has_image'], sanitize_text(image_info['format']), image_info['width'], 
                image_info['height'], image_info['file_size'], sanitize_text(image_info['url']), sanitize_text(image_info['base64']),
                sanitize_text(chart_info['type']), sanitize_text(chart_info['title']), sanitize_text(chart_info['data']), sanitize_text(chart_info['categories']), sanitize_text(chart_info['series']),
                hyperlink, shape_index, hidden,
                effects_info['shadow'], effects_info['glow'], effects_info['reflection'],
                effects_info['3d_effects'], sanitize_text(placeholder_type), sanitize_text(animation_effects)
            ]
            
            ws.append(row_data)
    
    return wb

# ------------------------------------------------
# Main
# ------------------------------------------------
def main():
    token = acquire_token_device_code()

    # Download PowerPoint
    ppt_url = f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content"
    print("⬇️ Downloading PowerPoint from OneDrive...")
    ppt_bytes = graph_get(ppt_url, token).content

    prs = Presentation(BytesIO(ppt_bytes))
    wb = extract_ppt_to_excel(prs)

    # Save Excel locally
    bio = BytesIO()
    wb.save(bio)
    bio.seek(0)

    # Upload Excel to OneDrive
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    dest_path = f"{DEST_FOLDER_ONEDRIVE}/Enhanced_PPT_Analysis_{timestamp}.xlsx"
    upload_url = f"{GRAPH_ROOT}{dest_path}:/content"

    print(f"📤 Uploading enhanced Excel analysis → {dest_path}")
    graph_put(upload_url, token, bio.getvalue(), content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")

    print("✅ Comprehensive PowerPoint analysis extracted and uploaded successfully!")
    print("📊 Extracted data includes:")
    print("   📍 Position & Size: EMU values, inches, rotation")
    print("   🎨 Font Details: name, size, bold, italic, underline, color, alignment")
    print("   📝 Text Formatting: line spacing, paragraph spacing")
    print("   🎨 Fill Properties: color, type (solid/gradient/pattern), transparency")
    print("   🖼️ Line Properties: color, width, style (dash/dot/solid)")
    print("   🖼️ Enhanced Image Info: format, dimensions, file size, URLs, base64")
    print("   � Chart Data: type, title, series, categories, data values")
    print("   �🔗 Hyperlinks and Z-order positioning")
    print("   ✨ Visual Effects: shadow, glow, reflection, 3D effects")
    print("   📋 Placeholder types and animation detection")
    print("   🔍 Shape visibility and advanced properties")
    print("   📊 Complete table data extraction")
    print("")
    print("📁 Images extracted to: ./extracted_images/ folder")
    print("🔗 Image URLs and base64 data included in Excel")
    print("📈 Chart data including series, categories, and values extracted")
    print("💡 Note: Install 'pillow' package for enhanced image dimension analysis:")
    print("   pip install pillow")

if __name__ == "__main__":
    main()
