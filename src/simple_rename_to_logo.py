import os
import requests
from io import BytesIO
from datetime import datetime
from dotenv import load_dotenv
from msal import PublicClientApplication
from pptx import Presentation

# ------------------------------------------------
# Load environment
# ------------------------------------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "")
CLIENT_ID = os.getenv("CLIENT_ID", "")
AUTHORITY = os.getenv("AUTHORITY", f"https://login.microsoftonline.com/{TENANT_ID}")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# ------------------------------------------------
# Auth helper
# ------------------------------------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]

# ------------------------------------------------
# Graph helpers
# ------------------------------------------------
def graph_get(url, token):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers, stream=True)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r

def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r

# ------------------------------------------------
# Simple renaming function
# ------------------------------------------------
def rename_google_shape_to_logo(prs):
    """Rename Google Shape;89;p16 to logo"""
    print("🎯 Looking for 'Google Shape;89;p16' to rename to 'logo'...")
    
    renamed_count = 0
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.name == "Google Shape;89;p16":
                old_name = shape.name
                shape.name = "logo"
                print(f"✅ Found and renamed on Slide {slide_num}:")
                print(f"   '{old_name}' → 'logo'")
                renamed_count += 1
    
    if renamed_count == 0:
        print("❌ Shape 'Google Shape;89;p16' not found")
        print("📋 Available shapes on Slide 1:")
        if len(prs.slides) > 0:
            for i, shape in enumerate(prs.slides[0].shapes):
                print(f"   {i+1}. '{shape.name}'")
        return False
    else:
        print(f"\n🎉 Successfully renamed {renamed_count} shape(s) to 'logo'!")
        return True

# ------------------------------------------------
# Main
# ------------------------------------------------
def main():
    print("🏷️  SIMPLE SHAPE RENAMER: Google Shape;89;p16 → logo")
    print("=" * 60)
    
    token = acquire_token_device_code()
    
    # Download PowerPoint
    print("⬇️ Downloading PowerPoint template...")
    try:
        ppt_response = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token)
        prs = Presentation(BytesIO(ppt_response.content))
    except Exception as e:
        print(f"❌ Failed to download: {e}")
        return
    
    # Attempt renaming
    success = rename_google_shape_to_logo(prs)
    
    if success:
        # Save and upload
        try:
            bio = BytesIO()
            prs.save(bio)
            bio.seek(0)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            dest_path = f"{DEST_FOLDER_ONEDRIVE}/Template_with_Logo_{timestamp}.pptx"
            upload_url = f"{GRAPH_ROOT}{dest_path}:/content"
            
            print(f"\n📤 Uploading renamed presentation → {dest_path}")
            graph_put(upload_url, token, bio.getvalue())
            print("✅ Successfully uploaded presentation with renamed shape!")
            print(f"   Shape 'Google Shape;89;p16' is now named 'logo' on Slide 1")
            
        except Exception as e:
            print(f"❌ Failed to upload: {e}")
            print("💾 Saving locally as backup...")
            with open("template_with_logo.pptx", "wb") as f:
                bio.seek(0)
                f.write(bio.read())
            print("✅ Saved locally as 'template_with_logo.pptx'")
    else:
        print("\n❌ No changes made - shape not found")

if __name__ == "__main__":
    main()