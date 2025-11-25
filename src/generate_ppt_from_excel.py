# import os
# import sys
# import json
# import requests
# from msal import PublicClientApplication
# from dotenv import load_dotenv
# from openpyxl import load_workbook

# # ---------------------
# # Load .env
# # ---------------------
# load_dotenv()
# print("🔍 Environment Variables Debug:")
# for var in ["TENANT_ID", "CLIENT_ID", "AUTHORITY", "LOCAL_EXCEL_PATH", "EXCEL_TITLE_CELL", "TEMPLATE_ONEDRIVE_PATH", "DEST_FOLDER_ONEDRIVE"]:
#     print(f"{var}: '{os.getenv(var, '')}'")

# TENANT_ID = os.getenv("TENANT_ID", "").strip()
# CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
# AUTHORITY = os.getenv("AUTHORITY", "").strip()
# LOCAL_EXCEL_PATH = os.getenv("LOCAL_EXCEL_PATH", "Book.xlsx")
# EXCEL_TITLE_CELL = os.getenv("EXCEL_TITLE_CELL", "A1")
# TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
# DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation.pptx")

# if not CLIENT_ID or not AUTHORITY:
#     print("❌ Missing required environment variables.")
#     sys.exit(1)

# # Use pre-consented Graph permissions
# SCOPES = ["https://graph.microsoft.com/.default"]
# GRAPH_ROOT = "https://graph.microsoft.com/v1.0"
# ME_DRIVE_ROOT = f"{GRAPH_ROOT}/me/drive/root:"

# # ---------------------
# # Read title from Excel
# # ---------------------
# def read_title_from_excel(path, cell_address):
#     try:
#         wb = load_workbook(path, data_only=True)
#         sheet = wb.active
#         value = sheet[cell_address].value
#         wb.close()
#         title = (value or "Untitled Presentation").strip()
#         print(f"📘 Title from Excel {path} {cell_address} => {title}")
#         return title
#     except Exception as e:
#         print(f"⚠️ Could not read Excel ({path}). Using default title. Error: {e}")
#         return "Untitled Presentation"

# # ---------------------
# # Token acquisition
# # ---------------------
# def acquire_token_device_code():
#     app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)

#     # Try cached token
#     accounts = app.get_accounts()
#     if accounts:
#         result = app.acquire_token_silent(SCOPES, account=accounts[0])
#         if result and "access_token" in result:
#             print("🔓 Using cached token.")
#             return result["access_token"]

#     # Device code flow
#     flow = app.initiate_device_flow(scopes=SCOPES)
#     if "user_code" not in flow:
#         raise RuntimeError("Failed to create device code flow. Check app config.")
#     print("🔑 To sign in, visit:", flow["verification_uri"])
#     print("🔑 Enter code:", flow["user_code"])
#     result = app.acquire_token_by_device_flow(flow)
#     if "access_token" not in result:
#         raise RuntimeError(f"Auth failed: {result.get('error_description') or result}")
#     return result["access_token"]

# # ---------------------
# # Graph helpers
# # ---------------------
# def graph_get(url, token, stream=False):
#     headers = {"Authorization": f"Bearer {token}"}
#     r = requests.get(url, headers=headers, stream=stream)
#     if r.status_code >= 400:
#         raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
#     return r

# def graph_put(url, token, data, content_type="application/octet-stream"):
#     headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
#     r = requests.put(url, headers=headers, data=data)
#     if r.status_code >= 400:
#         raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
#     return r

# # ---------------------
# # Main function
# # ---------------------
# def main():
#     title = read_title_from_excel(LOCAL_EXCEL_PATH, EXCEL_TITLE_CELL)
#     print(f"📘 Using title: {title}")

#     token = acquire_token_device_code()

#     # Download master deck
#     template_content_url = f"{ME_DRIVE_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content"
#     print(f"📥 Downloading template from {TEMPLATE_ONEDRIVE_PATH} ...")
#     r = graph_get(template_content_url, token, stream=True)
#     content = r.content
#     print(f"✅ Template size: {len(content)} bytes")

#     # Upload new presentation
#     dest_path = f"{DEST_FOLDER_ONEDRIVE}/{title}.pptx"
#     upload_url = f"{ME_DRIVE_ROOT}{dest_path}:/content"
#     print(f"📤 Uploading to {dest_path} ...")
#     graph_put(upload_url, token, data=content)
#     print("✅ Done. Check OneDrive:", dest_path)

# # ---------------------
# if __name__ == "__main__":
#     main()
# import os
# import sys
# import json
# import requests
# from io import BytesIO
# from msal import PublicClientApplication
# from dotenv import load_dotenv
# from openpyxl import load_workbook
# from pptx import Presentation

# # ---------------------
# # Load .env
# # ---------------------
# load_dotenv()
# print("🔍 Environment Variables Debug:")
# for var in [
#     "TENANT_ID",
#     "CLIENT_ID",
#     "AUTHORITY",
#     "EXCEL_ONEDRIVE_PATH",
#     "EXCEL_TITLE_CELL",
#     "TEMPLATE_ONEDRIVE_PATH",
#     "DEST_FOLDER_ONEDRIVE",
# ]:
#     print(f"{var}: '{os.getenv(var, '')}'")

# TENANT_ID = os.getenv("TENANT_ID", "").strip()
# CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
# AUTHORITY = os.getenv("AUTHORITY", "").strip() or f"https://login.microsoftonline.com/{TENANT_ID}"

# EXCEL_ONEDRIVE_PATH = os.getenv("EXCEL_ONEDRIVE_PATH", "/me/drive/root:/Book.xlsx")
# EXCEL_TITLE_CELL = os.getenv("EXCEL_TITLE_CELL", "A1")
# TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
# DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

# if not CLIENT_ID or not AUTHORITY:
#     print("❌ Missing required environment variables.")
#     sys.exit(1)

# SCOPES = ["https://graph.microsoft.com/.default"]
# GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# # ---------------------
# # Token acquisition
# # ---------------------
# def acquire_token_device_code():
#     app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)

#     accounts = app.get_accounts()
#     if accounts:
#         result = app.acquire_token_silent(SCOPES, account=accounts[0])
#         if result and "access_token" in result:
#             print("🔓 Using cached token.")
#             return result["access_token"]

#     flow = app.initiate_device_flow(scopes=SCOPES)
#     if "user_code" not in flow:
#         raise RuntimeError("Failed to create device code flow. Check Azure App configuration.")

#     print("🔑 To sign in, visit:", flow["verification_uri"])
#     print("🔑 Enter code:", flow["user_code"])
#     result = app.acquire_token_by_device_flow(flow)
#     if "access_token" not in result:
#         raise RuntimeError(f"Auth failed: {result.get('error_description') or result}")
#     return result["access_token"]

# # ---------------------
# # Graph helpers
# # ---------------------
# def graph_get(url, token, stream=False):
#     headers = {"Authorization": f"Bearer {token}"}
#     r = requests.get(url, headers=headers, stream=stream)
#     if r.status_code >= 400:
#         raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
#     return r

# def graph_put(url, token, data, content_type="application/octet-stream"):
#     headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
#     r = requests.put(url, headers=headers, data=data)
#     if r.status_code >= 400:
#         raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
#     return r

# def ensure_folder_exists(path, token):
#     """Create destination folder if missing"""
#     path_clean = path.split(":/")[-1].strip("/")
#     parts = path_clean.split("/")
#     base = "/me/drive/root:"
#     for i, part in enumerate(parts):
#         check_url = f"{GRAPH_ROOT}{base}/{part}"
#         resp = requests.get(check_url, headers={"Authorization": f"Bearer {token}"})
#         if resp.status_code == 404:
#             print(f"📁 Creating folder: {part}")
#             parent_path = base if i == 0 else f"{base}/{'/'.join(parts[:i])}"
#             create_url = f"{GRAPH_ROOT}{parent_path}:/children"
#             requests.post(
#                 create_url,
#                 headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
#                 json={"name": part, "folder": {}, "@microsoft.graph.conflictBehavior": "rename"},
#             )
#         base += f"/{part}"

# # ---------------------
# # Read Excel
# # ---------------------
# def read_excel_from_onedrive(token):
#     file_url = f"{GRAPH_ROOT}{EXCEL_ONEDRIVE_PATH}:/content"
#     print(f"📥 Downloading Excel from {EXCEL_ONEDRIVE_PATH} ...")
#     r = graph_get(file_url, token, stream=True)
#     wb = load_workbook(filename=BytesIO(r.content), data_only=True)
#     sheet = wb.active

#     headers = [cell.value for cell in sheet[1] if cell.value]
#     values = [cell.value for cell in sheet[2] if cell.value]

#     data = dict(zip(headers, values))
#     title = data.get(headers[0], "Untitled Presentation").strip()
#     wb.close()
#     print(f"📘 Title read from Excel ({headers[0]}) => {title}")
#     return title, data

# # ---------------------
# # Replace text recursively
# # ---------------------
# def replace_text_recursive(shape, replacements):
#     """Recursively replace text in grouped and normal shapes"""
#     if shape.shape_type == 6:  # GroupShape
#         for subshape in shape.shapes:
#             replace_text_recursive(subshape, replacements)
#     elif hasattr(shape, "text_frame") and shape.text_frame:
#         for paragraph in shape.text_frame.paragraphs:
#             for run in paragraph.runs:
#                 for key, value in replacements.items():
#                     placeholder = f"{{{{{key}}}}}"
#                     if placeholder in run.text:
#                         run.text = run.text.replace(placeholder, str(value))

# def apply_replacements(prs, replacements):
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             replace_text_recursive(shape, replacements)

# # ---------------------
# # PowerPoint Creation
# # ---------------------
# def create_ppt_from_template(token, title, replacements):
#     print(f"📥 Downloading template from {TEMPLATE_ONEDRIVE_PATH} ...")
#     template_url = f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content"
#     r = graph_get(template_url, token, stream=True)
#     prs = Presentation(BytesIO(r.content))

#     print("🧠 Applying replacements ...")
#     apply_replacements(prs, replacements)

#     ppt_stream = BytesIO()
#     prs.save(ppt_stream)
#     ppt_stream.seek(0)

#     ensure_folder_exists(DEST_FOLDER_ONEDRIVE, token)
#     dest_path = f"{DEST_FOLDER_ONEDRIVE}/{title}.pptx"
#     upload_url = f"{GRAPH_ROOT}{dest_path}:/content"

#     print(f"📤 Uploading to {dest_path} ...")
#     graph_put(upload_url, token, ppt_stream)
#     print("✅ Done. Uploaded successfully to OneDrive.")

# # ---------------------
# # Main
# # ---------------------
# def main():
#     token = acquire_token_device_code()
#     title, replacements = read_excel_from_onedrive(token)
#     create_ppt_from_template(token, title, replacements)

# if __name__ == "__main__":
#     main()
# import os
# import sys
# import requests
# from io import BytesIO
# from msal import PublicClientApplication
# from dotenv import load_dotenv
# from openpyxl import load_workbook
# from pptx import Presentation
# from pptx.dml.color import RGBColor

# # ---------------------
# # Load environment
# # ---------------------
# load_dotenv()

# TENANT_ID = os.getenv("TENANT_ID", "").strip()
# CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
# AUTHORITY = os.getenv("AUTHORITY", "").strip() or f"https://login.microsoftonline.com/{TENANT_ID}"

# EXCEL_ONEDRIVE_PATH = os.getenv("EXCEL_ONEDRIVE_PATH", "/me/drive/root:/Book.xlsx")
# TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
# DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

# SCOPES = ["https://graph.microsoft.com/.default"]
# GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# # ---------------------
# # Acquire token
# # ---------------------
# def acquire_token_device_code():
#     app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
#     accounts = app.get_accounts()
#     if accounts:
#         result = app.acquire_token_silent(SCOPES, account=accounts[0])
#         if result and "access_token" in result:
#             print("🔓 Using cached token.")
#             return result["access_token"]

#     flow = app.initiate_device_flow(scopes=SCOPES)
#     if "user_code" not in flow:
#         raise RuntimeError("Failed to create device code flow.")
#     print("🔑 Visit:", flow["verification_uri"])
#     print("🔑 Enter code:", flow["user_code"])
#     result = app.acquire_token_by_device_flow(flow)
#     if "access_token" not in result:
#         raise RuntimeError(f"Auth failed: {result}")
#     return result["access_token"]

# # ---------------------
# # Graph helpers
# # ---------------------
# def graph_get(url, token, stream=False):
#     headers = {"Authorization": f"Bearer {token}"}
#     r = requests.get(url, headers=headers, stream=stream)
#     if r.status_code >= 400:
#         raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
#     return r

# def graph_put(url, token, data, content_type="application/octet-stream"):
#     headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
#     r = requests.put(url, headers=headers, data=data)
#     if r.status_code >= 400:
#         raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
#     return r

# # ---------------------
# # Read Excel values
# # ---------------------
# def read_excel_values(token):
#     file_url = f"{GRAPH_ROOT}{EXCEL_ONEDRIVE_PATH}:/content"
#     print(f"📥 Downloading Excel from {EXCEL_ONEDRIVE_PATH} ...")
#     r = graph_get(file_url, token, stream=True)

#     wb = load_workbook(filename=BytesIO(r.content), data_only=True)
#     sheet = wb.active

#     data = {
#         "OverviewText": sheet["A2"].value or "",
#         "Header1": sheet["B2"].value or "",
#         "Description1": sheet["C2"].value or "",
#         "Header2": sheet["D2"].value or "",
#         "Description2": sheet["E2"].value or "",
#         "Header3": sheet["F2"].value or "",
#         "Description3": sheet["G2"].value or "",
#         "Header4": sheet["H2"].value or "",
#         "Description4": sheet["I2"].value or "",
#         "ThemeColor": sheet["J2"].value or "#3667B2",
#     }

#     wb.close()
#     print("📘 Excel Data Loaded:")
#     for k, v in data.items():
#         print(f"   {k}: {v}")
#     return data

# # ---------------------
# # Update PowerPoint Template
# # ---------------------
# def update_ppt_template(template_bytes, values):
#     prs = Presentation(BytesIO(template_bytes))

#     color_hex = values["ThemeColor"].lstrip("#")
#     color_rgb = RGBColor(int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16))

#     # Replace placeholders and color header/description text only
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text = shape.text
#                 for key, val in values.items():
#                     if key != "ThemeColor":
#                         text = text.replace(f"{{{{{key}}}}}", str(val))
#                 shape.text = text

#                 # Apply color only to headers and descriptions
#                 if any(
#                     key in shape.text
#                     for key in [
#                         values["Header1"], values["Header2"], values["Header3"], values["Header4"],
#                         values["Description1"], values["Description2"], values["Description3"], values["Description4"]
#                     ]
#                 ):
#                     for paragraph in shape.text_frame.paragraphs:
#                         for run in paragraph.runs:
#                             run.font.color.rgb = color_rgb

#     # Save updated presentation
#     bio = BytesIO()
#     prs.save(bio)
#     bio.seek(0)
#     print("🎨 Updated text placeholders and applied color to headers/descriptions only.")
#     return bio

# # ---------------------
# # Main
# # ---------------------
# def main():
#     token = acquire_token_device_code()
#     values = read_excel_values(token)

#     print(f"📥 Downloading template from {TEMPLATE_ONEDRIVE_PATH} ...")
#     r = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token, stream=True)
#     template_bytes = r.content
#     print(f"✅ Template size: {len(template_bytes)} bytes")

#     # Modify PPT
#     updated_ppt = update_ppt_template(template_bytes, values)

#     # Upload new presentation
#     dest_path = f"{DEST_FOLDER_ONEDRIVE}/{values['OverviewText'] or 'Updated_Presentation'}.pptx"
#     upload_url = f"{GRAPH_ROOT}{dest_path}:/content"
#     print(f"📤 Uploading updated presentation to {dest_path} ...")
#     graph_put(upload_url, token, data=updated_ppt.getvalue())
#     print("✅ Done. Check OneDrive:", dest_path)


# if __name__ == "__main__":
#     main()
import os
import sys
import requests
from io import BytesIO
from msal import PublicClientApplication
from dotenv import load_dotenv
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor

# ---------------------
# Load environment
# ---------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "").strip()
CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
AUTHORITY = os.getenv("AUTHORITY", "").strip() or f"https://login.microsoftonline.com/{TENANT_ID}"

EXCEL_ONEDRIVE_PATH = os.getenv("EXCEL_ONEDRIVE_PATH", "/me/drive/root:/Book.xlsx")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# ---------------------
# Acquire token
# ---------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    if "user_code" not in flow:
        raise RuntimeError("Failed to create device code flow.")
    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]

# ---------------------
# Graph helpers
# ---------------------
def graph_get(url, token, stream=False):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers, stream=stream)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r

def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r

# ---------------------
# Read Excel values
# ---------------------
def read_excel_values(token):
    file_url = f"{GRAPH_ROOT}{EXCEL_ONEDRIVE_PATH}:/content"
    print(f"📥 Downloading Excel from {EXCEL_ONEDRIVE_PATH} ...")
    r = graph_get(file_url, token, stream=True)

    wb = load_workbook(filename=BytesIO(r.content), data_only=True)
    sheet = wb.active

    data = {
        "OverviewText": sheet["A2"].value or "",
        "Header1": sheet["B2"].value or "",
        "Description1": sheet["C2"].value or "",
        "Header2": sheet["D2"].value or "",
        "Description2": sheet["E2"].value or "",
        "Header3": sheet["F2"].value or "",
        "Description3": sheet["G2"].value or "",
        "Header4": sheet["H2"].value or "",
        "Description4": sheet["I2"].value or "",
        "ThemeColor": sheet["K2"].value or "#3667B2",  # S100 color
        "SlideTitle": sheet["L2"].value or "Generated Slide",  # Slide title
    }

    wb.close()
    print("📘 Excel Data Loaded:")
    for k, v in data.items():
        print(f"   {k}: {v}")
    return data

# ---------------------
# Update PowerPoint Template
# ---------------------
def update_ppt_template(template_bytes, values):
    prs = Presentation(BytesIO(template_bytes))

    # 🎨 Convert HEX color
    color_hex = values["ThemeColor"].lstrip("#")
    color_rgb = RGBColor(int(color_hex[0:2], 16),
                         int(color_hex[2:4], 16),
                         int(color_hex[4:6], 16))

    # Background shapes to recolor
    bg_shapes = {"OverviewText_bg", "OverviewText_bg_1", "OverviewText_bg_2"}

    for slide in prs.slides:
        for shape in slide.shapes:

            # 🎨 Apply dynamic background color
            if shape.name in bg_shapes:
                try:
                    fill = shape.fill
                    fill.solid()
                    fill.fore_color.rgb = color_rgb
                    print(f"🎨 Recolored {shape.name} → #{color_hex.upper()}")
                except Exception as e:
                    print(f"⚠️ Could not recolor {shape.name}: {e}")

            # 🧠 Replace placeholders (including SlideTitle)
            if shape.has_text_frame:
                text = shape.text
                for key, val in values.items():
                    if key != "ThemeColor":
                        text = text.replace(f"{{{{{key}}}}}", str(val))
                shape.text = text

                # Apply color to text dynamically
                if any(key in shape.text for key in [
                    values["Header1"], values["Header2"], values["Header3"], values["Header4"],
                    values["Description1"], values["Description2"], values["Description3"], values["Description4"],
                    values["SlideTitle"]
                ]):
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = color_rgb

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    print("✅ Updated text placeholders, title, and background color.")
    return bio

# ---------------------
# Main
# ---------------------
def main():
    token = acquire_token_device_code()
    values = read_excel_values(token)

    print(f"📥 Downloading template from {TEMPLATE_ONEDRIVE_PATH} ...")
    r = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token, stream=True)
    template_bytes = r.content
    print(f"✅ Template size: {len(template_bytes)} bytes")

    # Modify PPT
    updated_ppt = update_ppt_template(template_bytes, values)

    # Safe filename from SlideTitle
    safe_title = "".join(c for c in values["SlideTitle"] if c.isalnum() or c in (" ", "_", "-")).strip()
    if not safe_title:
        safe_title = "Generated_Presentation"

    dest_path = f"{DEST_FOLDER_ONEDRIVE}/{safe_title.replace(' ', '_')}.pptx"

    upload_url = f"{GRAPH_ROOT}{dest_path}:/content"
    print(f"📤 Uploading updated presentation to {dest_path} ...")
    graph_put(upload_url, token, data=updated_ppt.getvalue())
    print("✅ Done. Check OneDrive:", dest_path)


if __name__ == "__main__":
    main()
