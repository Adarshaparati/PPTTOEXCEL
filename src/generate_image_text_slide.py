import os
import requests
from io import BytesIO
from msal import PublicClientApplication
from dotenv import load_dotenv
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor

# ---------------------
# Load environment
# ---------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "").strip()
CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
AUTHORITY = os.getenv("AUTHORITY", "").strip() or f"https://login.microsoftonline.com/{TENANT_ID}"

EXCEL_ONEDRIVE_PATH = os.getenv("EXCEL_ONEDRIVE_PATH", "/me/drive/root:/Book.xlsx")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# ---------------------
# Token acquisition
# ---------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    if "user_code" not in flow:
        raise RuntimeError("Failed to create device code flow.")

    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]

# ---------------------
# Graph helpers
# ---------------------
def graph_get(url, token, stream=False):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers, stream=stream)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r

def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r

# ---------------------
# Read Excel configuration
# ---------------------
def read_excel_values(token):
    file_url = f"{GRAPH_ROOT}{EXCEL_ONEDRIVE_PATH}:/content"
    print(f"📥 Downloading Excel from {EXCEL_ONEDRIVE_PATH} ...")
    r = graph_get(file_url, token, stream=True)

    wb = load_workbook(filename=BytesIO(r.content), data_only=True)
    sheet = wb["Sheet2"]  # Adjust sheet name if needed

    data = {
        "Slide_No": int(sheet["A2"].value or 1),
        "Slide_Title": sheet["B2"].value or "Title Here",
        "Slide_Text": sheet["C2"].value or "Description Here",
        "Image_Path": sheet["D2"].value or "",
        "P100": sheet["E2"].value or "#3667B2",
        "S100": sheet["F2"].value or "#000000",
    }

    wb.close()
    print("📘 Excel Data Loaded:")
    for k, v in data.items():
        print(f"   {k}: {v}")
    return data

# ---------------------
# Update PowerPoint Template
# ---------------------
def update_ppt_template(template_bytes, values):
    prs = Presentation(BytesIO(template_bytes))
    slide_index = values["Slide_No"] - 1  # 0-indexed

    if slide_index >= len(prs.slides):
        raise IndexError(f"Slide {values['Slide_No']} not found in template.")

    slide = prs.slides[slide_index]

    # Convert HEX → RGBColor
    def hex_to_rgb(hex_str):
        hex_str = hex_str.lstrip("#")
        return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

    bg_color = hex_to_rgb(values["P100"])
    accent_color = hex_to_rgb(values["S100"])

    # 🖼️ Download image
    image_bytes = None
    if values["Image_Path"]:
        try:
            print(f"🖼️ Downloading image: {values['Image_Path']}")
            img_res = requests.get(values["Image_Path"])
            img_res.raise_for_status()
            image_bytes = BytesIO(img_res.content)
            print("✅ Image downloaded successfully.")
        except Exception as e:
            print(f"⚠️ Failed to download image: {e}")

    # 🔄 Update placeholders and backgrounds
    for shape in slide.shapes:
        # Replace title and text
        if shape.has_text_frame:
            text = shape.text
            text = text.replace("{{Title}}", str(values["Slide_Title"]))
            text = text.replace("{{Description}}", str(values["Slide_Text"]))
            shape.text = text

            # Apply accent color to text
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = accent_color

        # Replace only the main image (named "Image")
        elif shape.name == "Image" and image_bytes:
            try:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(image_bytes, left, top, width, height)
                print(f"🖼️ Replaced image in shape '{shape.name}'")
            except Exception as e:
                print(f"⚠️ Could not replace image: {e}")

        # Apply S100 to Image_bg_1 and Image_bg_2
        elif shape.name in ("Image_bg_1", "Image_bg_2"):
            try:
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = accent_color
                print(f"🎨 Recolored {shape.name} → {values['S100']}")
            except Exception as e:
                print(f"⚠️ Could not recolor {shape.name}: {e}")

        # Apply P100 to any text background shapes
        elif "Text_bg" in shape.name:
            try:
                fill = shape.fill
                fill.solid()
                fill.fore_color.rgb = bg_color
                print(f"🎨 Recolored {shape.name} → {values['P100']}")
            except Exception as e:
                print(f"⚠️ Could not recolor {shape.name}: {e}")

    # Save updated PPT
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    print("✅ Slide updated successfully.")
    return bio

# ---------------------
# Main
# ---------------------
def main():
    token = acquire_token_device_code()
    values = read_excel_values(token)

    print(f"📥 Downloading template from {TEMPLATE_ONEDRIVE_PATH} ...")
    r = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token, stream=True)
    template_bytes = r.content
    print(f"✅ Template size: {len(template_bytes)} bytes")

    updated_ppt = update_ppt_template(template_bytes, values)

    safe_title = "".join(c for c in values["Slide_Title"] if c.isalnum() or c in (" ", "_", "-")).strip()
    dest_path = f"{DEST_FOLDER_ONEDRIVE}/{safe_title.replace(' ', '_')}_ImageText.pptx"
    upload_url = f"{GRAPH_ROOT}{dest_path}:/content"

    print(f"📤 Uploading updated presentation to {dest_path} ...")
    graph_put(upload_url, token, data=updated_ppt.getvalue())
    print("✅ Done. Check OneDrive:", dest_path)


if __name__ == "__main__":
    main()
