import os
import requests
from io import BytesIO
from datetime import datetime
from msal import PublicClientApplication
from dotenv import load_dotenv
from openpyxl import load_workbook
from pptx import Presentation
from pptx.dml.color import RGBColor

# ------------------------------------------------
# Load environment
# ------------------------------------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "").strip()
CLIENT_ID = os.getenv("CLIENT_ID", "").strip()
AUTHORITY = os.getenv("AUTHORITY", "").strip() or f"https://login.microsoftonline.com/{TENANT_ID}"

EXCEL_ONEDRIVE_PATH = os.getenv("EXCEL_ONEDRIVE_PATH", "/me/drive/root:/Book.xlsx")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"


# ------------------------------------------------
# Authentication
# ------------------------------------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]


# ------------------------------------------------
# Graph helpers
# ------------------------------------------------
def graph_get(url, token, stream=False):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers, stream=stream)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r


def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r


# ------------------------------------------------
# Helpers
# ------------------------------------------------
def hex_to_rgb(hex_str):
    """Convert #RRGGBB string to RGBColor safely."""
    if not hex_str:
        return RGBColor(0, 0, 0)

    hex_str = str(hex_str).strip().replace("#", "")
    if len(hex_str) != 6 or any(c not in "0123456789ABCDEFabcdef" for c in hex_str):
        print(f"⚠️ Invalid hex color '{hex_str}', defaulting to black")
        return RGBColor(0, 0, 0)

    return RGBColor(int(hex_str[0:2], 16),
                    int(hex_str[2:4], 16),
                    int(hex_str[4:6], 16))


# ------------------------------------------------
# 1️⃣ TEXT SLIDE (Sheet1)
# ------------------------------------------------
def update_text_slide(prs, sheet):
    print("🧾 Updating TEXT SLIDE (Sheet1)...")

    # Safely extract slide number from column M
    try:
        slide_no = int(sheet["M2"].value)
    except (ValueError, TypeError):
        slide_no = 8
        print(f"⚠️ Invalid slide number found in M2 — defaulting to slide {slide_no}")

    # Collect slide data from Excel columns
    data = {
        "OverviewText": sheet["A2"].value or "",
        "Header1": sheet["B2"].value or "",
        "Description1": sheet["C2"].value or "",
        "Header2": sheet["D2"].value or "",
        "Description2": sheet["E2"].value or "",
        "Header3": sheet["F2"].value or "",
        "Description3": sheet["G2"].value or "",
        "Header4": sheet["H2"].value or "",
        "Description4": sheet["I2"].value or "",
        "P100": sheet["J2"].value or "#3667B2",
        "S100": sheet["K2"].value or "#8A8B8C",
        "SlideTitle": sheet["L2"].value or "Default Title",
        "Slide_No": slide_no,
        "Icon1": sheet["N2"].value or "",
        "Icon2": sheet["O2"].value or "",
        "Icon3": sheet["P2"].value or "",
        "Icon4": sheet["Q2"].value or "",
    }

    # Convert hex colors safely
    bg_color = hex_to_rgb(data["P100"])
    accent_color = hex_to_rgb(data["S100"])

    # Get the slide
    try:
        slide = prs.slides[data["Slide_No"] - 1]
    except IndexError:
        print(f"❌ Slide {data['Slide_No']} not found in template.")
        return

    # Replace text placeholders
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        text = shape.text
        replacements = {
            "{{OverviewText}}": data["OverviewText"],
            "{{Header1}}": data["Header1"],
            "{{Description1}}": data["Description1"],
            "{{Header2}}": data["Header2"],
            "{{Description2}}": data["Description2"],
            "{{Header3}}": data["Header3"],
            "{{Description3}}": data["Description3"],
            "{{Header4}}": data["Header4"],
            "{{Description4}}": data["Description4"],
            "{{SlideTitle}}": data["SlideTitle"],
        }

        for key, value in replacements.items():
            text = text.replace(key, str(value))

        shape.text = text

    # Apply colors to shapes
    for shape in slide.shapes:
        shape_name = shape.name.lower().strip()
        try:
            if "overviewtext_bg_1" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = accent_color  # S100 color
            elif "overviewtext_bg" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg_color  # P100 color
            elif "bg_primary" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg_color  # P100 color
            elif "bg_secondary" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = accent_color
        except Exception as e:
            print(f"⚠️ Could not color shape '{shape.name}': {e}")

    # Replace icons
    icon_mapping = {
        "icon1": data["Icon1"],
        "icon2": data["Icon2"],
        "icon3": data["Icon3"],
        "icon4": data["Icon4"],
    }

    for shape in slide.shapes:
        shape_name = shape.name.lower().strip()
        for icon_name, icon_url in icon_mapping.items():
            if icon_name in shape_name and icon_url:
                try:
                    # Download icon image
                    icon_res = requests.get(icon_url)
                    icon_res.raise_for_status()
                    icon_bytes = BytesIO(icon_res.content)
                    
                    # Replace the image
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    slide.shapes._spTree.remove(shape._element)
                    slide.shapes.add_picture(icon_bytes, left, top, width, height)
                    print(f"🎨 Icon replaced in '{shape.name}'.")
                    break
                except Exception as e:
                    print(f"⚠️ Could not replace icon '{shape.name}': {e}")

    print(f"✅ Text slide {data['Slide_No']} updated successfully.")


# ------------------------------------------------
# 2️⃣ TEXT + IMAGE SLIDE (Sheet2)
# ------------------------------------------------
def update_text_image_slide(prs, sheet):
    print("\n🖼️ Updating TEXT + IMAGE SLIDE (Sheet2)...")

    data = {
        "Slide_No": int(sheet["A2"].value or 9),
        "Slide_Title": sheet["B2"].value or "Text Image Slide",
        "Slide_Text": sheet["C2"].value or "Description here",
        "Image_Path": sheet["D2"].value or "",
        "P100": sheet["E2"].value or "#3667B2",
        "S100": sheet["F2"].value or "#000000",
    }

    bg_color = hex_to_rgb(data["P100"])
    accent_color = hex_to_rgb(data["S100"])

    try:
        slide = prs.slides[data["Slide_No"] - 1]
    except IndexError:
        print(f"❌ Slide {data['Slide_No']} not found in template.")
        return

    # Download image
    image_bytes = None
    if data["Image_Path"]:
        try:
            img_res = requests.get(data["Image_Path"])
            img_res.raise_for_status()
            image_bytes = BytesIO(img_res.content)
            print("🖼️ Image downloaded successfully.")
        except Exception as e:
            print(f"⚠️ Image download failed: {e}")

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text.replace("{{Title}}", data["Slide_Title"])
            text = text.replace("{{Description}}", data["Slide_Text"])
            shape.text = text

            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    r.font.color.rgb = accent_color

        elif shape.name.lower().strip() == "image" and image_bytes:
            try:
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                slide.shapes._spTree.remove(shape._element)
                slide.shapes.add_picture(image_bytes, left, top, width, height)
                print(f"🖼️ Image replaced in '{shape.name}'.")
            except Exception as e:
                print(f"⚠️ Could not replace image: {e}")

    # Apply colors to shapes
    for shape in slide.shapes:
        shape_name = shape.name.lower().strip()
        try:
            if "image_bg_1" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = accent_color  # S100 color
            elif "image_bg_2" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = accent_color  # S100 color
            elif "image_bg_3" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = bg_color  # P100 color
        except Exception as e:
            print(f"⚠️ Could not color shape '{shape.name}': {e}")

    print("✅ Text+Image slide updated.")


# ------------------------------------------------
# 3️⃣ TABLE SLIDE (Sheet3)
# ------------------------------------------------
def update_table_slide(prs, sheet):
    print("\n📊 Updating TABLE SLIDE (Sheet3)...")

    # Base data
    data = {
        "Slide_No": int(sheet["A2"].value or 3),
        "Title": sheet["B2"].value or "Data and Analysis",
        "RowHeader1": sheet["C2"].value or "Sales Data",
        "RowHeader2": sheet["D2"].value or "Customer Data",
        "RowHeader3": sheet["E2"].value or "Market Trends",
        "RowHeader4": sheet["F2"].value or "Comparisons",
        "Column_Header1": sheet["G2"].value or "Value 1",
        "Column_Header2": sheet["H2"].value or "Value 2",
        "Column_Header3": sheet["I2"].value or "Value 3",
        "Column_Header4": sheet["J2"].value or "Value 4",
        "P100": sheet["K2"].value or "#3667B2",
        "S100": sheet["L2"].value or "#8A8B8C",
    }

    # Explicit mapping for all 16 cell values (M → AB)
    col_letters = [
        "M", "N", "O", "P",  # Row 1
        "Q", "R", "S", "T",  # Row 2
        "U", "V", "W", "X",  # Row 3
        "Y", "Z", "AA", "AB" # Row 4
    ]

    # Assign dynamically
    cell_keys = [
        "C1R1", "C2R1", "C3R1", "C4R1",
        "C1R2", "C2R2", "C3R2", "C4R2",
        "C1R3", "C2R3", "C3R3", "C4R3",
        "C1R4", "C2R4", "C3R4", "C4R4",
    ]

    for key, col in zip(cell_keys, col_letters):
        data[key] = sheet[f"{col}2"].value or ""

    # Convert colors
    p100_color = hex_to_rgb(data["P100"])
    s100_color = hex_to_rgb(data["S100"])

    # Get the slide
    try:
        slide = prs.slides[data["Slide_No"] - 1]
    except IndexError:
        print(f"❌ Slide {data['Slide_No']} not found in template.")
        return

    # Replace placeholders
    for shape in slide.shapes:
        if shape.has_text_frame:
            original_text = shape.text
            text = shape.text

            replacements = {
                "{{Title}}": str(data["Title"]),
                "{{RowHeader1}}": str(data["RowHeader1"]),
                "{{RowHeader2}}": str(data["RowHeader2"]),
                "{{RowHeader3}}": str(data["RowHeader3"]),
                "{{RowHeader4}}": str(data["RowHeader4"]),
                "{{Column_Header1}}": str(data["Column_Header1"]),
                "{{Column_Header2}}": str(data["Column_Header2"]),
                "{{Column_Header3}}": str(data["Column_Header3"]),
                "{{Column_Header4}}": str(data["Column_Header4"]),
            }

            # Add value placeholders (C1R1_VALUE → Excel C1R1)
            for key in cell_keys:
                replacements[f"{{{{{key}_VALUE}}}}"] = str(data[key])

            # Apply replacements
            for k, v in replacements.items():
                text = text.replace(k, v)
            shape.text = text

            # Apply text color based on shape name
            shape_name = shape.name.lower().strip()
            
            # Check if this shape contains C1R values (column 1 values) or Column_Header1 BEFORE replacement
            if any(f"{{{{C1R{i}_VALUE}}}}" in original_text for i in range(1, 5)) or "{{Column_Header1}}" in original_text:
                # Apply white color to C1R values and Column_Header1
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = RGBColor(255, 255, 255)  # White
            else:
                # Apply default text color
                for p in shape.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.color.rgb = s100_color

        # Apply color fills
        shape_name = shape.name.lower().strip()
        try:
            if "column1" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = s100_color  # S100 color for COLUMN1
            elif "column_header" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = p100_color
            elif "rowheader" in shape_name:
                shape.fill.solid()
                shape.fill.fore_color.rgb = s100_color
        except:
            pass

    print("✅ Table slide (4x4) updated successfully with values from Excel.")

# ------------------------------------------------
# MAIN
# ------------------------------------------------
def main():
    token = acquire_token_device_code()

    # Download Excel and PPT template
    excel_bytes = graph_get(f"{GRAPH_ROOT}{EXCEL_ONEDRIVE_PATH}:/content", token, stream=True).content
    wb = load_workbook(filename=BytesIO(excel_bytes), data_only=True)

    r = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token, stream=True)
    prs = Presentation(BytesIO(r.content))

    # Update slides
    update_text_slide(prs, wb["Sheet1"])
    update_text_image_slide(prs, wb["Sheet2"])
    update_table_slide(prs, wb["Sheet3"])

    # Save and upload
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)

    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    dest_path = f"{DEST_FOLDER_ONEDRIVE}/Full_Auto_Presentation_{timestamp}.pptx"
    upload_url = f"{GRAPH_ROOT}{dest_path}:/content"

    print(f"\n📤 Uploading final combined presentation → {dest_path}")
    graph_put(upload_url, token, data=bio.getvalue())
    print("✅ All 3 slides generated and uploaded successfully!")


if __name__ == "__main__":
    main()
