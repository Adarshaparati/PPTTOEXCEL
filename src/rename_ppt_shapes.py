import os
import requests
from io import BytesIO
from datetime import datetime
from dotenv import load_dotenv
from msal import PublicClientApplication
from pptx import Presentation

# ------------------------------------------------
# Load environment
# ------------------------------------------------
load_dotenv()

TENANT_ID = os.getenv("TENANT_ID", "")
CLIENT_ID = os.getenv("CLIENT_ID", "")
AUTHORITY = os.getenv("AUTHORITY", f"https://login.microsoftonline.com/{TENANT_ID}")
TEMPLATE_ONEDRIVE_PATH = os.getenv("TEMPLATE_ONEDRIVE_PATH", "/me/drive/root:/IntroductionTemplate.pptx")
DEST_FOLDER_ONEDRIVE = os.getenv("DEST_FOLDER_ONEDRIVE", "/me/drive/root:/Presentation")

SCOPES = ["https://graph.microsoft.com/.default"]
GRAPH_ROOT = "https://graph.microsoft.com/v1.0"

# ------------------------------------------------
# Auth helper
# ------------------------------------------------
def acquire_token_device_code():
    app = PublicClientApplication(CLIENT_ID, authority=AUTHORITY)
    accounts = app.get_accounts()
    if accounts:
        result = app.acquire_token_silent(SCOPES, account=accounts[0])
        if result and "access_token" in result:
            print("🔓 Using cached token.")
            return result["access_token"]

    flow = app.initiate_device_flow(scopes=SCOPES)
    print("🔑 Visit:", flow["verification_uri"])
    print("🔑 Enter code:", flow["user_code"])
    result = app.acquire_token_by_device_flow(flow)
    if "access_token" not in result:
        raise RuntimeError(f"Auth failed: {result}")
    return result["access_token"]

# ------------------------------------------------
# Graph helpers
# ------------------------------------------------
def graph_get(url, token):
    headers = {"Authorization": f"Bearer {token}"}
    r = requests.get(url, headers=headers, stream=True)
    if r.status_code >= 400:
        raise RuntimeError(f"GET {url} failed: {r.status_code} {r.text}")
    return r

def graph_put(url, token, data, content_type="application/octet-stream"):
    headers = {"Authorization": f"Bearer {token}", "Content-Type": content_type}
    r = requests.put(url, headers=headers, data=data)
    if r.status_code >= 400:
        raise RuntimeError(f"PUT {url} failed: {r.status_code} {r.text}")
    return r

# ------------------------------------------------
# Shape renaming functions
# ------------------------------------------------

def list_all_shapes(prs):
    """List all shapes in the presentation with their current names and slide numbers"""
    print("\n📋 CURRENT SHAPE INVENTORY:")
    print("=" * 70)
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n🎯 SLIDE {slide_num}:")
        print("-" * 40)
        
        for i, shape in enumerate(slide.shapes):
            shape_type = "Unknown"
            content_preview = ""
            
            # Determine shape type and content
            if hasattr(shape, 'shape_type'):
                type_mapping = {
                    1: "AutoShape", 2: "Callout", 3: "Chart", 4: "Comment", 
                    5: "Freeform", 6: "Group", 7: "Line", 8: "LinkedOLEObject",
                    9: "LinkedPicture", 10: "Media", 11: "OLEObject", 12: "Picture",
                    13: "Placeholder", 14: "TextBox", 15: "3DModel", 16: "Canvas",
                    17: "Connector", 18: "Ink", 19: "Table", 20: "SmartArt"
                }
                shape_type = type_mapping.get(shape.shape_type, f"Type_{shape.shape_type}")
            
            # Get content preview
            if shape.has_text_frame and shape.text.strip():
                content_preview = shape.text.strip()[:50] + ("..." if len(shape.text.strip()) > 50 else "")
            elif shape.shape_type == 12:  # Picture
                content_preview = "[IMAGE]"
            elif shape.shape_type == 19:  # Table
                content_preview = "[TABLE]"
            
            print(f"   {i+1:2d}. '{shape.name}' ({shape_type})")
            if content_preview:
                print(f"       Content: {content_preview}")

def rename_shapes_by_rules(prs):
    """Rename shapes based on predefined rules and patterns"""
    print("\n🏷️  APPLYING SMART RENAMING RULES...")
    print("=" * 50)
    
    renaming_stats = {"renamed": 0, "skipped": 0}
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n🎯 Processing Slide {slide_num}...")
        
        # Counters for each slide
        counters = {
            "text": 1, "image": 1, "table": 1, "icon": 1, 
            "background": 1, "button": 1, "shape": 1
        }
        
        for shape in slide.shapes:
            old_name = shape.name
            new_name = None
            
            # Rule 1: Text content analysis
            if shape.has_text_frame and shape.text.strip():
                text_content = shape.text.strip().lower()
                
                # Check for specific text patterns
                if any(word in text_content for word in ["title", "heading", "header"]):
                    new_name = f"Title_{counters['text']}"
                    counters['text'] += 1
                elif any(word in text_content for word in ["description", "overview", "content"]):
                    new_name = f"Description_{counters['text']}"
                    counters['text'] += 1
                elif "{{" in text_content and "}}" in text_content:
                    # Placeholder text
                    placeholder_text = text_content.replace("{{", "").replace("}}", "")
                    new_name = f"Placeholder_{placeholder_text[:20]}"
                else:
                    new_name = f"Text_{counters['text']}"
                    counters['text'] += 1
            
            # Rule 2: Shape type analysis
            elif hasattr(shape, 'shape_type'):
                if shape.shape_type == 12:  # Picture
                    new_name = f"Image_{counters['image']}"
                    counters['image'] += 1
                elif shape.shape_type == 19:  # Table
                    new_name = f"Table_{counters['table']}"
                    counters['table'] += 1
                elif shape.shape_type == 1:  # AutoShape (could be background)
                    # Check if it's likely a background based on size
                    if shape.width > 5000000 and shape.height > 3000000:  # Large shapes likely backgrounds
                        new_name = f"Background_{counters['background']}"
                        counters['background'] += 1
                    else:
                        new_name = f"Shape_{counters['shape']}"
                        counters['shape'] += 1
                else:
                    new_name = f"Element_{counters['shape']}"
                    counters['shape'] += 1
            
            # Rule 3: Position-based naming (small shapes could be icons)
            if new_name and shape.width < 1000000 and shape.height < 1000000:  # Small shapes
                if "Image" in new_name:
                    new_name = f"Icon_{counters['icon']}"
                    counters['icon'] += 1
            
            # Apply the new name
            if new_name and new_name != old_name:
                shape.name = new_name
                print(f"   ✅ Renamed: '{old_name}' → '{new_name}'")
                renaming_stats["renamed"] += 1
            else:
                print(f"   ⏭️  Kept: '{old_name}'")
                renaming_stats["skipped"] += 1
    
    print(f"\n📊 RENAMING SUMMARY:")
    print(f"   ✅ Renamed: {renaming_stats['renamed']} shapes")
    print(f"   ⏭️  Skipped: {renaming_stats['skipped']} shapes")

def rename_shapes_interactive(prs):
    """Interactive mode to rename shapes manually"""
    print("\n🖱️  INTERACTIVE RENAMING MODE")
    print("=" * 40)
    print("Commands: 'q' to quit, 's' to skip shape, 'l' to list all")
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n🎯 SLIDE {slide_num}:")
        
        for i, shape in enumerate(slide.shapes):
            print(f"\n   Shape {i+1}: '{shape.name}'")
            
            # Show shape details
            if shape.has_text_frame and shape.text.strip():
                print(f"   Content: {shape.text.strip()[:100]}")
            
            # Get user input
            while True:
                new_name = input(f"   New name (current: '{shape.name}'): ").strip()
                
                if new_name.lower() == 'q':
                    return
                elif new_name.lower() == 's':
                    break
                elif new_name.lower() == 'l':
                    list_all_shapes(prs)
                    continue
                elif new_name and new_name != shape.name:
                    shape.name = new_name
                    print(f"   ✅ Renamed to: '{new_name}'")
                    break
                elif not new_name:
                    break

def rename_shapes_from_mapping(prs, mapping_dict):
    """Rename shapes based on a provided mapping dictionary"""
    print("\n📝 APPLYING CUSTOM MAPPING...")
    print("=" * 40)
    
    renamed_count = 0
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            old_name = shape.name
            if old_name in mapping_dict:
                new_name = mapping_dict[old_name]
                shape.name = new_name
                print(f"   Slide {slide_num}: '{old_name}' → '{new_name}'")
                renamed_count += 1
    
    print(f"\n✅ Renamed {renamed_count} shapes using custom mapping")

def rename_specific_shape(prs, target_name, new_name):
    """Rename a specific shape by its exact name"""
    print(f"\n🎯 RENAMING SPECIFIC SHAPE...")
    print("=" * 40)
    print(f"Looking for: '{target_name}'")
    print(f"Renaming to: '{new_name}'")
    
    found = False
    
    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if shape.name == target_name:
                old_name = shape.name
                shape.name = new_name
                print(f"✅ Found and renamed on Slide {slide_num}:")
                print(f"   '{old_name}' → '{new_name}'")
                found = True
                return True
    
    if not found:
        print(f"❌ Shape '{target_name}' not found in presentation")
        print("\n📋 Available shapes:")
        for slide_num, slide in enumerate(prs.slides, 1):
            print(f"   Slide {slide_num}:")
            for shape in slide.shapes:
                print(f"     - '{shape.name}'")
    
    return found

# ------------------------------------------------
# Predefined renaming templates
# ------------------------------------------------

def apply_standard_template_names(prs):
    """Apply standard naming convention for common presentation elements"""
    print("\n🎨 APPLYING STANDARD TEMPLATE NAMING...")
    print("=" * 45)
    
    # Common naming patterns for presentation elements
    standard_mappings = {
        # Text elements
        "title": "SlideTitle",
        "heading": "Header1", 
        "subtitle": "Subtitle",
        
        # Background elements
        "background": "Background_Primary",
        "bg": "Background_Secondary",
        
        # Content areas
        "content": "ContentArea",
        "text": "TextContent",
        "description": "Description",
        
        # Media elements
        "image": "MainImage",
        "picture": "MainImage",
        "icon": "Icon1",
        
        # Navigation/UI
        "button": "ActionButton",
        "nav": "Navigation"
    }
    
    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n🎯 Slide {slide_num}:")
        
        for shape in slide.shapes:
            old_name = shape.name.lower()
            new_name = None
            
            # Check for matches in standard mappings
            for pattern, standard_name in standard_mappings.items():
                if pattern in old_name:
                    new_name = standard_name
                    break
            
            # Apply naming based on content and position
            if not new_name and shape.has_text_frame:
                text = shape.text.strip().lower()
                if any(word in text for word in ["overview", "description"]):
                    new_name = "OverviewText"
                elif "header" in text or "title" in text:
                    new_name = "SlideTitle"
            
            if new_name and new_name != shape.name:
                shape.name = new_name
                print(f"   ✅ '{shape.name}' → '{new_name}'")

# ------------------------------------------------
# Main execution
# ------------------------------------------------

def main():
    token = acquire_token_device_code()
    
    # Download PowerPoint
    print("⬇️ Downloading PowerPoint template...")
    ppt_response = graph_get(f"{GRAPH_ROOT}{TEMPLATE_ONEDRIVE_PATH}:/content", token)
    prs = Presentation(BytesIO(ppt_response.content))
    
    print("\n🎭 POWERPOINT SHAPE RENAMING TOOL")
    print("=" * 50)
    
    while True:
        print("\n📋 AVAILABLE OPTIONS:")
        print("1. 📋 List all current shape names")
        print("2. 🎯 Rename specific shape (Google Shape;72;p16 → logo)")
        print("3. 🤖 Auto-rename using smart rules") 
        print("4. 🎨 Apply standard template naming")
        print("5. 🖱️  Interactive renaming mode")
        print("6. 💾 Save and upload modified presentation")
        print("7. 🚪 Exit without saving")
        
        choice = input("\nSelect option (1-7): ").strip()
        
        if choice == "1":
            list_all_shapes(prs)
        
        elif choice == "2":
            # Rename Google Shape;72;p16 to logo
            target_shape = "Google Shape;72;p16"
            new_name = "logo"
            success = rename_specific_shape(prs, target_shape, new_name)
            if success:
                print(f"\n🎉 Successfully renamed '{target_shape}' to '{new_name}'!")
        
        elif choice == "3":
            rename_shapes_by_rules(prs)
        
        elif choice == "4":
            apply_standard_template_names(prs)
        
        elif choice == "5":
            rename_shapes_interactive(prs)
        
        elif choice == "6":
            # Save and upload
            bio = BytesIO()
            prs.save(bio)
            bio.seek(0)
            
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            dest_path = f"{DEST_FOLDER_ONEDRIVE}/Renamed_Template_{timestamp}.pptx"
            upload_url = f"{GRAPH_ROOT}{dest_path}:/content"
            
            print(f"\n📤 Uploading renamed presentation → {dest_path}")
            graph_put(upload_url, token, bio.getvalue())
            print("✅ Presentation with renamed shapes uploaded successfully!")
            break
        
        elif choice == "7":
            print("👋 Exiting without saving...")
            break
        
        else:
            print("❌ Invalid option. Please choose 1-7.")

if __name__ == "__main__":
    main()